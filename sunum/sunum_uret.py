# -*- coding: utf-8 -*-
"""sunum_uret.py — Tanıtım sunumunu tek kaynaktan üretir.

═══════════════════════════════════════════════════════════════════════════════
Sekiz dosyayı da ÜRETİR; hiçbiri elle düzenlenmez:

    Akilli_Konaklama_Tanitim.pptx        ← Türkçe sunum (ekran)
    Akilli_Konaklama_Tanitim.pdf         ← PowerPoint ile dışa aktarılır
    Akilli_Konaklama_Tanitim.html        ← slayt PNG'leri gömülü mobil sunum
    Akilli_Konaklama_Tanitim_Baski.pptx  ← beyaz zemin (mono yazıcı)
    Akilli_Konaklama_Tanitim_Baski.pdf
    Akilli_Konaklama_Intro_EN.pptx       ← İngilizce sunum ve karşılıkları
    Akilli_Konaklama_Intro_EN.pdf
    Akilli_Konaklama_Intro_EN.html
    Akilli_Konaklama_Intro_EN_Print.pptx / .pdf

Neden betik? Çünkü elle düzenlemek dört ayrı iş demek
-----------------------------------------------------
Sunum tek bir dosya değil: iki dil × (ekran + baskı) × (pptx + pdf + html).
Bir cümleyi elle değiştirmek bunların hepsine tek tek dokunmayı gerektirir ve
ilk atlanan dosyada sunum kendi kendisiyle çelişmeye başlar. İçerik burada
BİR kez yazılır; geri kalan her şey türetilir.

Neden şablon kopyalamak yerine sıfırdan kuruluyor
--------------------------------------------------
Depoda tasarımı taşıyan bir "kaynak sunum" yok; olsaydı ikili bir dosyayı
sürüm kontrolünde tutmak ve her düzenlemeden sonra elle güncellemek
gerekirdi. Bunun yerine slaytların geometrisi ve renkleri bu dosyada
SAYIYLA duruyor (bkz. YERLEŞİM ve PALET). Kazanç: baskı sürümü ayrı bir
dönüştürme adımı değil, aynı yapının ikinci bir paletle yeniden kurulması.
Renkleri sonradan XML üzerinde eşlemeye çalışmak, rengi tanımsız bırakılmış
şekillerde sessizce yanlış sonuç verir.

Sayılar nereden geliyor
------------------------
Sürüm ``pyproject.toml``'dan OKUNUR — sunumda elle güncellenen sürüm numarası
kalmasın diye. Ölçüm sayıları (test, kapsam, izin, bulgu) ``docs/TEST_REPORT.md``
ve ``docs/SECURITY_REVIEW.md`` içindeki ÇALIŞTIRILMIŞ çıktılardan alınmıştır;
her birinin yanında hangi belgeden geldiği yazılıdır. Bu belgeler yeniden
ölçüldüğünde OLCUM bölümü de güncellenmelidir.

Kullanım:
    python sunum_uret.py                 # sekiz dosyanın tamamı
    python sunum_uret.py --sadece-pptx   # PDF/HTML/baskı atlanır
    python sunum_uret.py --dil tr        # tek dil
    python sunum_uret.py --kontrol       # üretmeden metin sığma denetimi

Gereksinim: python-pptx, pywin32 (PDF/PNG için PowerPoint). Çıkış kodu 0 = başarılı.
"""

import base64
import glob
import os
import re
import shutil
import sys
import tempfile

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:  # pragma: no cover - eski konsollar
    pass


def _iz(isaret):
    """Konsol UTF-8 değilse işareti ASCII karşılığına düşürür."""
    yedek = {"✓": "[OK]", "✗": "[!]", "→": "->", "·": "-"}
    try:
        isaret.encode(sys.stdout.encoding or "utf-8")
        return isaret
    except Exception:
        return yedek.get(isaret, "?")


KOK = os.path.dirname(os.path.abspath(__file__))
PROJE = os.path.dirname(KOK)
SIMGE = os.path.join(PROJE, "app", "ui", "resources", "icons", "app.png")

#: Uygulama ekran görüntüleri. ekran_yakala.py geçici bir demo veritabanı
#: kurup her sayfayı bu klasöre PNG olarak basar; sunum onları gömer.
#: Görüntü eksikse üretim BAŞLARKEN durur — eksik görselle sessizce
#: üretilen sunum, ekranı boş bir çerçeveyle göstermekten kötüdür.
EKRAN_DIZIN = os.path.join(KOK, "ekranlar")

DILLER = ("tr", "en")

#: Dosya adları. Türkçe "Tanitim", İngilizce "Intro_EN" — aynı önek ikisini
#: klasörde yan yana tutar.
ADLAR = {
    "tr": ("Akilli_Konaklama_Tanitim", "Akilli_Konaklama_Tanitim_Baski"),
    "en": ("Akilli_Konaklama_Intro_EN", "Akilli_Konaklama_Intro_EN_Print"),
}


def _yollar(dil):
    """(pptx, pdf, html) — ekran sürümü."""
    ad = ADLAR[dil][0]
    return (os.path.join(KOK, ad + ".pptx"),
            os.path.join(KOK, ad + ".pdf"),
            os.path.join(KOK, ad + ".html"))


def _baski_yollari(dil):
    """(pptx, pdf) — beyaz zeminli baskı sürümü."""
    ad = ADLAR[dil][1]
    return os.path.join(KOK, ad + ".pptx"), os.path.join(KOK, ad + ".pdf")


def _surum():
    """Sürümü pyproject.toml'dan okur — sunumda elle güncellenen sayı kalmasın."""
    try:
        with open(os.path.join(PROJE, "pyproject.toml"), encoding="utf-8") as f:
            for satir in f:
                eslesme = re.match(r'^version\s*=\s*"([^"]+)"', satir.strip())
                if eslesme:
                    return eslesme.group(1)
    except OSError:
        pass
    return "0.1.0"


#: Ölçüm tarihi — üretim anında `datetime.date.today()` ile doldurulur.
#: Elle yazılan bir tarih, ölçümler yenilendiğinde sessizce eskiyordu.
def _bugun():
    from datetime import date

    bugun = date.today()
    aylar_tr = ("Ocak", "Şubat", "Mart", "Nisan", "Mayıs", "Haziran", "Temmuz",
                "Ağustos", "Eylül", "Ekim", "Kasım", "Aralık")
    aylar_en = ("January", "February", "March", "April", "May", "June", "July",
                "August", "September", "October", "November", "December")
    return {
        "tr": "%d %s %d" % (bugun.day, aylar_tr[bugun.month - 1], bugun.year),
        "en": "%d %s %d" % (bugun.day, aylar_en[bugun.month - 1], bugun.year),
    }


OLCUM_TARIHI = _bugun()

#: Slaytlardaki sayılar KAYNAK KODDAN ölçülür (bkz. sunum/olcum.py).
#:
#: Neden? Önceki sürümde test adedi, kapsam yüzdesi ve tablo sayısı slayt
#: metnine elle yazılmıştı ve yanlarında "bu sayılar ölçüldü" diyen bir
#: dipnot vardı. Kod değiştiğinde sayılar değişmiyor, sunum sessizce
#: yanlışlaşıyordu. Artık her sayı üretim anında hesaplanır; hesaplanamayan
#: bir sayı slayttan ÇIKARILIR, tahmin edilmez.
_OLCUM_ONBELLEK = {}


def olcumler(hizli=False):
    """Ölçüm sözlüğünü döndürür (bir kez hesaplanır)."""
    anahtar = bool(hizli)
    if anahtar not in _OLCUM_ONBELLEK:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        try:
            from olcum import olc

            _OLCUM_ONBELLEK[anahtar] = olc(testleri_calistir=not hizli)
        except Exception as hata:  # pragma: no cover - ölçüm aracı yoksa
            print("   ! olcum modulu calistirilamadi: %s" % hata)
            _OLCUM_ONBELLEK[anahtar] = {}
    return _OLCUM_ONBELLEK[anahtar]
YAYIN_AYI = {"tr": "Ağustos 2026", "en": "August 2026"}
DEPO = "github.com/Azizsekerdil/smart-hotel-management-system"


# ══════════════════════════════════════════════════════════════════════════
# PALET
#
# İki palet, tek yapı. Ekran paleti koyu (#0B0F14); baskı paleti beyaz zemin
# ve KOYULAŞTIRILMIŞ vurgu renkleri kullanır — çünkü 18 slaydın her birinin
# arka planı tam sayfa koyu renktir ve mono lazerde sayfa başına ~%95 toner
# demektir.
#
# Baskı renkleri iki ölçüte göre seçildi:
#   1) Beyaz zeminde kontrast ≥ 4.5 (WCAG AA, gövde metni)
#   2) Kart zeminleri beyaz DEĞİL, çok açık gri — sınır görünmezse sayfa
#      duvara döner; #F5F7F9 mono lazerde ~%7 toner.
#
# İçerik renklere ROL adıyla başvurur ("marka", "altin", ...). Böylece iki
# palet arasında birebir karşılık zorunlu olur: bir rol eklendiğinde
# diğerinde de tanımlanmazsa üretim KeyError ile durur, sessizce yanlış
# renkte basmaz.
# ══════════════════════════════════════════════════════════════════════════

PALET_EKRAN = {
    "zemin": "0B0F14",
    "kart": "141920",
    "cerceve": "2B3442",
    "yazi": "E6EDF3",
    "soluk": "8B98A9",
    "marka": "4A9EE0",
    "gok": "38BDF8",
    "teal": "2DD4BF",
    "altin": "F6D365",
    "yesil": "22C55E",
    "mor": "A78BFA",
    "kirmizi": "F87171",
    # İkon dairelerinin zemini: emoji kendi renklerini taşıdığı için daire
    # yalnızca arkalıktır, koyu kalmalıdır.
    "daire1": "1B5E9E",
    "daire2": "0F8F8A",
    "daire3": "5A3D7A",
    "daire4": "8A5A12",
}

PALET_BASKI = {
    "zemin": "FFFFFF",
    "kart": "F5F7F9",
    "cerceve": "DDE3E9",
    "yazi": "10171F",
    "soluk": "4B5563",
    "marka": "145A8F",
    "gok": "1A5F8F",
    "teal": "0B6F68",
    "altin": "8A6410",
    "yesil": "15803D",
    "mor": "5B21B6",
    "kirmizi": "8C1515",
    "daire1": "E3F0FA",
    "daire2": "E3EFEE",
    "daire3": "EDE7F3",
    "daire4": "FBF3DE",
}

#: Kart ikonlarının daire rengi sırayla döner; içerik renk seçmez, sıra seçer.
DAIRE_DONGUSU = ("daire1", "daire2", "daire3", "daire4")


# ══════════════════════════════════════════════════════════════════════════
# YERLEŞİM (inç) — 16:9, 13,333 × 7,5
#
# Değerler göz kararıyla değil, tek bir ızgaradan türetildi: sol/sağ kenar
# 0,60; başlık bandı 0,32–1,07; içerik 1,70–6,68; dipnot 6,80. Bir kutunun
# yerini değiştirmek isteyen buradaki sayıyı değiştirir — slayt kurucularında
# gömülü sabit yoktur.
# ══════════════════════════════════════════════════════════════════════════

SLAYT_G, SLAYT_Y = 13.3333, 7.5
ORTA = SLAYT_G / 2

BASLIK_KUTU = (0.60, 0.32, 12.10, 0.75)
ALT_KUTU = (0.66, 1.02, 12.00, 0.44)
DIPNOT_KUTU = (0.66, 6.80, 12.00, 0.40)

# Kart ızgarası (2×2). Kart yüksekliği en uzun gövde metnine göre seçildi:
# 13 punto Calibri 5,43 inçlik kutuda satır başına ~70 karakter alır, bütçe
# 275 karakter, yani en fazla dört satır (~0,87 inç). Daha yüksek bir kart
# metnin altında kocaman bir boşluk bırakıyordu (ölçüldü: ilk üretimde
# kartların yarısı boştu).
KART_X = (0.60, 6.80)
KART_T = (1.74, 4.32)
KART_G, KART_Y = 5.95, 2.15
IKON_OFS, IKON_CAP = (0.22, 0.20), 0.52
AD_OFS, AD_OLCU = (0.88, 0.20), (4.90, 0.52)
GOVDE_OFS, GOVDE_OLCU = (0.26, 0.84), (5.43, 1.18)

# Sayı ızgarası (4×2) + geniş not kutusu
SAYI_X = (0.60, 3.75, 6.90, 10.05)
SAYI_T = (1.75, 3.60)
SAYI_G, SAYI_Y = 2.90, 1.55
RAKAM_OFS, RAKAM_OLCU = (0.00, 0.14), (2.90, 0.80)
ACIKLAMA_OFS, ACIKLAMA_OLCU = (0.10, 0.95), (2.70, 0.52)
NOT_KUTU = (0.60, 5.55, 12.35, 1.30)

# Kapak / kapanış
PUL_X = (1.00, 3.88, 6.76, 9.64)
PUL_OLCU = (2.72, 0.55)

# Ekran görüntüsü slaytları. Görüntüler 16:9'dur (1600×900); tek görselde
# yükseklik sabitlenir, genişlik oranından türetilir ve yatay ortalanır.
# Çift görselde kart ızgarasının sütunları (KART_X/KART_G) kullanılır —
# görsel genişliği sabitlenir, yükseklik oranından gelir (5,95 × 9/16 ≈ 3,35).
EKRAN_TEK_UST, EKRAN_TEK_YUK = 1.58, 5.10
EKRAN_CIFT_ETIKET_T, EKRAN_CIFT_T = 2.56, 3.00

# Punto
P_BASLIK_IKON, P_BASLIK = 28, 30
P_ALT, P_DIPNOT = 13, 11
P_KART_AD, P_KART_GOVDE, P_KART_IKON = 15, 13, 17
P_RAKAM, P_ACIKLAMA = 32, 11
P_KAPAK_BASLIK, P_KAPAK_ALT, P_KAPAK_SLOGAN, P_PUL = 40, 20, 17, 12

YAZI_TIPI = "Calibri"

#: Metin sığma bütçeleri (karakter). --kontrol bunlarla ölçer. Sayılar
#: kutu genişliği ÷ ortalama karakter genişliği × satır sayısıyla bulundu;
#: aşan metin kutudan taşar ve slaytta üst üste biner.
BUTCE = {
    "baslik": 62,
    "alt": 210,
    "dipnot": 205,
    "kart_ad": 34,
    "kart_govde": 275,
    "aciklama": 92,
    "kapak_baslik": 40,
    "kapak_alt": 110,
    "pul": 26,
    "not": 260,
    "ekran_etiket": 34,
}


# ══════════════════════════════════════════════════════════════════════════
# İÇERİK
#
# Her metin iki dilde yazılır: ("türkçe", "english"). Çeviri sonradan yapılan
# bir adım değil, içeriğin kendisidir — böylece bir slayt eklenip İngilizcesi
# unutulamaz: metin ikili değilse üretim başlarken durur (bkz. _M).
#
# Slayt türleri:
#   kapak  — büyük başlık, slogan, dört pul, alt bilgi
#   sayi   — 4×2 sayı kutusu + geniş not
#   kart   — başlık + alt başlık + 2×2 kart + dipnot
#   ekran  — başlık + alt başlık + 1 büyük ya da 2 yan yana ekran görüntüsü
#            (görseller sunum/ekranlar/ içinden; ekran_yakala.py üretir)
# ══════════════════════════════════════════════════════════════════════════


def _M(deger, dil):
    """İki dilli metinden istenen dili verir; ikili değilse hata verir."""
    if isinstance(deger, tuple) and len(deger) == 2:
        return deger[0 if dil == "tr" else 1]
    raise TypeError("Metin iki dilli olmalı: (tr, en) — alınan: %r" % (deger,))


SLAYTLAR = [
    # ── 1. Kapak ─────────────────────────────────────────────────────────
    {
        "tip": "kapak",
        "simge": True,
        "pul": True,
        "baslik": ("AKILLI KONAKLAMA YÖNETİM SİSTEMİ",
                   "SMART HOSPITALITY MANAGEMENT SYSTEM"),
        "alt": ("Otel · butik otel · pansiyon · apart otel · tatil köyü için "
                "yapay zekâ destekli Windows ve macOS masaüstü PMS",
                "AI-assisted Windows & macOS desktop PMS for hotels, boutique "
                "hotels, guest houses, aparthotels and resorts"),
        "slogan": ("Doğru Sat.  Doğru Ölç.  Veriyi Koru.",
                   "Sell Right.  Measure Right.  Protect the Data."),
        "puller": [
            ("Windows 10/11 + macOS", "Windows 10/11 & macOS"),
            ("%100 yerel yapay zekâ", "100% local AI option"),
            ("Şifreli kimlik verisi", "Encrypted identity data"),
            ("Türkçe / İngilizce", "Turkish / English"),
        ],
        "bilgi": ("Sürüm __SURUM__ · __AY__ · MIT lisansı · Yasal uyum sorumluluğu "
                  "işletmeye aittir",
                  "Version __SURUM__ · __AY__ · MIT license · Legal compliance "
                  "remains the operator's responsibility"),
    },

    # ── 2. Sayılarla ─────────────────────────────────────────────────────
    {
        "tip": "sayi",
        "baslik": ("🔢  Sayılarla Program", "🔢  The Program in Numbers"),
        "kutular": [
            ("__TABLO__", "marka",
             ("Veritabanı tablosu · __MODEL__ ORM modeli · Alembic ile sürümlenir",
              "Database tables · __MODEL__ ORM models · versioned with Alembic")),
            ("__TEST__", "teal",
             ("Toplanan test — tamamı ağa çıkmadan çalışır; bu kural ayrı bir testle korunur",
              "Collected tests — all run without network access; that rule is itself tested")),
            ("__KAPSAM__", "gok",
             ("Dal dahil kapsam — `pytest --cov=app --cov-branch` çıktısından",
              "Branch coverage — taken from `pytest --cov=app --cov-branch`")),
            ("__IZIN__", "altin",
             ("İzin · __ROL__ varsayılan rol · yetki iki katmanda birden denetlenir",
              "Permissions · __ROL__ default roles · checked in two separate layers")),
            ("13", "yesil",
             ("Kritik senaryo — çakışma, no-show, kota, bozuk JSON dahil",
              "Critical scenarios — overlap, no-show, quota, malformed JSON")),
            ("__SAGLAYICI__", "mor",
             ("Yapay zekâ sağlayıcısı: LM Studio · OpenAI · NVIDIA · Anthropic · Mock",
              "AI providers: LM Studio · OpenAI · NVIDIA · Anthropic · Mock")),
            ("__EKRAN__", "marka",
             ("Sunumdaki uygulama ekranı — hepsi uydurma demo veriyle üretildi",
              "Application screens in this deck — all captured on fictional demo data")),
            ("__BANDIT_YUKSEK__", "kirmizi",
             ("bandit'te yüksek önemde bulgu · orta: __BANDIT_ORTA__ · düşük: __BANDIT_DUSUK__",
              "High-severity bandit findings · medium: __BANDIT_ORTA__ · low: __BANDIT_DUSUK__")),
        ],
        "not": ("Bu sayıların hiçbiri slayta elle yazılmadı: sunum üretilirken "
                "`sunum/olcum.py` tarafından kaynak koddan ölçüldü (__TARIH__). "
                "Ölçülemeyen bir değerin yerine rakam değil \"ölçülmedi\" basılır. "
                "Yöntem: docs/TEST_REPORT.md ve docs/SECURITY_REVIEW.md.",
                "None of these numbers was typed in: `sunum/olcum.py` measures them "
                "from the source as the deck is built (__TARIH__). What cannot be "
                "measured prints as \"not measured\", never as a remembered figure."),
    },

    # ── 3. Kimin için ────────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🏨  Kimin İçin — ve Ne Değil",
                   "🏨  Who It Is For — and What It Is Not"),
        "alt": ("Program tek bir işletmenin bilgisayarında çalışır: bulut aboneliği "
                "değil, kanal yöneticisi değil, muhasebe paketi değil.",
                "The program runs on a single property's own computer: it is not a "
                "cloud subscription, not a channel manager, not an accounting suite."),
        "kartlar": [
            (("🏩", "Hedef İşletme", "Target Property"),
             ("5–150 odalı otel, butik otel, pansiyon, apart otel ve tatil köyü. "
              "Tek bilgisayarda ya da yerel ağda birkaç kullanıcıyla çalışır; "
              "internet bağlantısı zorunlu değildir.",
              "Hotels, boutique hotels, guest houses, aparthotels and resorts with "
              "5–150 rooms. Runs on one machine or a small local network; an "
              "internet connection is not required.")),
            (("🧾", "Kapsadığı İş", "What It Covers"),
             ("Rezervasyondan çıkışa kadar ön büro döngüsünün tamamı: satış, "
              "konaklama, folyo, tahsilat, kat hizmetleri, teknik servis ve "
              "yönetim raporları.",
              "The complete front-office cycle from booking to departure: sales, "
              "stay, folio, payment, housekeeping, maintenance and management "
              "reports.")),
            (("🚫", "Ne Değil", "What It Is Not"),
             ("Kanal yöneticisi, POS ve muhasebe programı değildir. e-Fatura ve "
              "Kimlik Bildirim Sistemi entegrasyonları YAPILMADI: veri modeli "
              "hazır, bağlantı yok.",
              "Not a channel manager, not a POS, not accounting software. Turkish "
              "e-invoice and police-notification integrations are NOT DONE: the "
              "data model exists, the connection does not.")),
            (("⚖️", "Sorumluluk Sınırı", "Where Responsibility Sits"),
             ("Yazılım veri işlemeyi kolaylaştırır; KVKK, e-Fatura ve Kimlik "
              "Bildirim yükümlülüklerine uyum sorumluluğu tamamen kullanıcı "
              "işletmeye aittir.",
              "The software makes data processing easier; compliance with data "
              "protection, e-invoicing and guest-notification obligations remains "
              "entirely with the operating business.")),
        ],
        "dipnot": ("Tamamlanmamış modüller varsayılan olarak kapalıdır ve arayüzde "
                   "durumu yazılı bir bekleme ekranıyla gösterilir — docs/ROADMAP.md",
                   "Unfinished modules are disabled by default and shown in the UI "
                   "as a placeholder screen that states their status — docs/ROADMAP.md"),
    },

    # ── 4. Günün akışı ───────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🔄  Bir Günün Akışı — Dört Ekran",
                   "🔄  A Day in Four Screens"),
        "alt": ("Gün panelde başlar, ön büroda geçer, kat ve teknik ekipleriyle "
                "sürer, gün sonu kapanışıyla biter.",
                "The day starts on the dashboard, runs through the front desk and "
                "the housekeeping and maintenance teams, and ends with night audit."),
        "kartlar": [
            (("📊", "Yönetim Paneli", "Dashboard"),
             ("Doluluk, ADR, RevPAR ve ALOS kartları; kritik uyarılar (geciken "
              "çıkış, açık arıza, kritik stok) ve 14 günlük doluluk eğrisi tek "
              "ekranda toplanır.",
              "Occupancy, ADR, RevPAR and ALOS tiles; critical alerts (overdue "
              "departures, open faults, low stock) and a 14-day occupancy curve, "
              "all on one screen.")),
            (("📅", "Rezervasyonlar", "Reservations"),
             ("Takvim, grup rezervasyonu, bekleme listesi, kanal bazlı kayıt. "
              "Gece gece fiyat hesabı ve çakışma engelleme kayıt anında çalışır.",
              "Calendar, group bookings, waiting list, channel-tagged records. "
              "Night-by-night pricing and overlap prevention run at save time.")),
            (("🛎", "Ön Büro", "Front Desk"),
             ("Giriş, oda kartı, refakatçi, folyo, ek ücret, tahsilat ve çıkış. "
              "Erken giriş ve geç çıkış ücretleri dilim dilim hesaplanır.",
              "Check-in, room card, companions, folio, extra charges, payment and "
              "check-out. Early check-in and late check-out are charged in slots.")),
            (("🧹", "Kat ve Teknik", "Housekeeping & Maintenance"),
             ("Günlük temizlik görevleri üretilir ve atanır; arıza kaydı "
              "gerektiğinde odayı satışa kapatır — kapalı oda rezervasyona düşmez.",
              "Daily cleaning tasks are generated and assigned; a fault record can "
              "take a room out of sale — a blocked room cannot be booked.")),
        ],
        "dipnot": ("Kısayollar: Ctrl+1 Panel · Ctrl+2 Rezervasyon · Ctrl+3 Ön Büro · "
                   "Ctrl+4 Odalar · Ctrl+5 Misafirler · Ctrl+R Raporlar",
                   "Shortcuts: Ctrl+1 Dashboard · Ctrl+2 Reservations · Ctrl+3 Front "
                   "Desk · Ctrl+4 Rooms · Ctrl+5 Guests · Ctrl+R Reports"),
    },

    # ── 5. Ekran: Yönetim Paneli ─────────────────────────────────────────
    {
        "tip": "ekran",
        "baslik": ("🖥️  Ekranda: Yönetim Paneli", "🖥️  On Screen: The Dashboard"),
        "alt": ("Doluluk, boş oda, günlük gelir, ADR ve RevPAR kartları; kritik "
                "uyarılar ve 14 günlük doluluk tahmini — güne tek ekrandan başlanır.",
                "Occupancy, vacant rooms, daily revenue, ADR and RevPAR tiles; "
                "critical alerts and a 14-day occupancy forecast — the day starts "
                "on a single screen."),
        "gorseller": [("01_dashboard.png", None)],
        "dipnot": ("Bu ve sonraki ekran görüntüleri programın kendisinden alınmıştır; "
                   "verilerin tamamı sentetik demo kaydıdır ve gerçek kişi verisi içermez.",
                   "This and the following screenshots are taken from the program "
                   "itself; all of the data is synthetic demo content with no real "
                   "personal data."),
    },

    # ── 6. Ekran: Ön Büro ────────────────────────────────────────────────
    {
        "tip": "ekran",
        "baslik": ("🛎️  Ekranda: Ön Büro", "🛎️  On Screen: The Front Desk"),
        "alt": ("Bekleyen giriş ve çıkışlar, oteldeki misafirler ve tahsil edilmemiş "
                "folyo bakiyesi tek ekranda; giriş işlemi listeden tek düğmeyle başlar.",
                "Pending arrivals and departures, in-house guests and the uncollected "
                "folio balance on one screen; check-in starts from the list with a "
                "single button."),
        "gorseller": [("03_frontdesk.png", None)],
    },

    # ── 7. Ekran: Kat Hizmetleri ve Teknik Servis ────────────────────────
    {
        "tip": "ekran",
        "baslik": ("🧹  Ekranda: Kat Hizmetleri ve Teknik Servis",
                   "🧹  On Screen: Housekeeping and Maintenance"),
        "alt": ("Günün temizlik görevleri tek düğmeyle üretilir, görevliye atanır ve "
                "kontrol adımıyla kapanır; arıza kayıtları öncelik ve oda ilişkisiyle izlenir.",
                "Daily cleaning tasks are generated with one click, assigned and closed "
                "with an inspection step; fault records are tracked with priority and "
                "room links."),
        "gorseller": [("06_housekeeping.png", ("Kat Hizmetleri", "Housekeeping")),
                      ("07_maintenance.png", ("Teknik Servis", "Maintenance"))],
    },

    # ── 8. Rezervasyon motoru ────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🧮  Aynı Oda Aynı Gece İki Kez Satılamaz",
                   "🧮  The Same Room Cannot Be Sold Twice"),
        "alt": ("Çakışma engelleme bir uyarı değil, iki katmanlı bir kısıttır: "
                "uygulama kontrolü ile veritabanı kısıtı aynı kuralı ayrı ayrı uygular.",
                "Overlap prevention is not a warning but a two-layer constraint: the "
                "application check and the database constraint enforce the same rule "
                "independently."),
        "kartlar": [
            (("🚧", "İki Aşamalı Kontrol", "Two-Stage Check"),
             ("Önce uygulama katmanı müsaitliği sorgular ve anlaşılır bir hata "
              "üretir; kayıt anında veritabanı kısıtı aynı kuralı yeniden uygular. "
              "Yarış durumunda ikinci işlem reddedilir.",
              "First the application layer queries availability and produces a "
              "readable error; at write time the database constraint applies the "
              "same rule again. In a race, the second transaction is rejected.")),
            (("📐", "Yarı-Açık Aralık", "Half-Open Interval"),
             ("Tarihler [giriş, çıkış) olarak modellenir: çıkış günü gece "
              "sayılmaz, oda aynı gün yeniden satılabilir. Çakışma hatalarının "
              "çoğu tam bu sınırda doğar.",
              "Dates are modelled as [arrival, departure): the departure day is not "
              "a night, so the room can be sold again the same day. Most overlap "
              "bugs are born exactly at this boundary.")),
            (("💵", "Gece Gece Fiyat", "Night-by-Night Pricing"),
             ("Sezon fiyatı, hafta sonu farkı, ekstra kişi, indirim ve vergi her "
              "gece için ayrı hesaplanır; tek bir ortalama fiyatla çarpma yapılmaz.",
              "Season rate, weekend uplift, extra person, discount and tax are "
              "computed for each night separately; no single average rate is "
              "multiplied by the stay length.")),
            (("🔀", "Durum Makinesi", "State Machine"),
             ("Taslak → onaylı → giriş → çıkış; iptal ve no-show ayrı yollardır. "
              "Geçersiz geçişler reddedilir ve ücret kuralları duruma bağlıdır.",
              "Draft → confirmed → checked-in → checked-out; cancellation and "
              "no-show are separate paths. Invalid transitions are rejected and fee "
              "rules depend on the state.")),
        ],
        "dipnot": ("99 alan testi yalnız bu katmanı sınar: müsaitlik 24, fiyat 23, "
                   "durum makinesi 18, para ve tarih değer nesneleri 34.",
                   "99 domain tests cover this layer alone: availability 24, pricing "
                   "23, state machine 18, money and date value objects 34."),
    },

    # ── 9. Ekran: Rezervasyonlar ─────────────────────────────────────────
    {
        "tip": "ekran",
        "baslik": ("📅  Ekranda: Rezervasyonlar", "📅  On Screen: Reservations"),
        "alt": ("Arama, durum ve tarih filtreleri; onay numarası, gece sayısı, tutar "
                "ve bakiye tek listede. İptal ve opsiyonlu kayıtlar renkle ayrışır.",
                "Search with status and date filters; confirmation number, nights, "
                "amount and balance in one list. Cancelled and optional bookings are "
                "colour-coded."),
        "gorseller": [("02_reservations.png", None)],
    },

    # ── 10. Ekran: Oda Planı ─────────────────────────────────────────────
    {
        "tip": "ekran",
        "baslik": ("🛏️  Ekranda: Oda Planı", "🛏️  On Screen: The Room Map"),
        "alt": ("Bina ve kat bazlı oda planı: doluluk ve temizlik durumu renk koduyla "
                "görünür; seçilen odanın ayrıntısı ve temizlik işlemleri sağ paneldedir.",
                "A building- and floor-based room map: occupancy and cleanliness are "
                "colour-coded; the selected room's details and actions sit in the "
                "right-hand panel."),
        "gorseller": [("04_rooms.png", None)],
    },

    # ── 11. Folyo ────────────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🧾  Folyo — Silinmez, Gerekçeyle Geçersiz Kılınır",
                   "🧾  Folio — Never Deleted, Only Voided with a Reason"),
        "alt": ("Mali kayıtta düzeltmek silmek değildir. Program ücret satırını "
                "silmez; gerekçesiyle birlikte geçersiz kılar ve denetim izi kalır.",
                "In financial records, correcting is not deleting. The program never "
                "removes a charge line; it voids it with a reason and the audit "
                "trail survives."),
        "kartlar": [
            (("✍️", "Denetim İzi Korunur", "The Trail Survives"),
             ("Yanlış girilen ücret satırı iptal edilir ve gerekçesi kayda geçer. "
              "Böylece gün sonu toplamı ile geçmiş her zaman açıklanabilir kalır.",
              "A mistaken charge line is voided and its reason recorded. The night "
              "audit total and the history therefore stay explainable at any later "
              "date.")),
            (("💳", "Tahsilat ve İade", "Payments and Refunds"),
             ("Nakit, kart ve havale tahsilatı; iade, depozito ve hasar kaydı. "
              "Bakiye folyodan hesaplanır; tutar tutarsızsa işlem reddedilir.",
              "Cash, card and transfer payments; refunds, deposits and damage "
              "records. The balance is derived from the folio; an inconsistent "
              "amount is rejected.")),
            (("⏱", "Erken Giriş / Geç Çıkış", "Early In / Late Out"),
             ("Ücret dilim dilim hesaplanır: 4 saatlik geç çıkış iki dilimdir ve "
              "gecelik fiyatın %50'sine karşılık gelir. Kural testle sabitlenmiştir.",
              "Charged in slots: a 4-hour late check-out is two slots and equals "
              "50% of the nightly rate. The rule is pinned down by a test.")),
            (("📆", "Gün Sonu Kapanışı", "Night Audit"),
             ("Girişler, çıkışlar, konaklayanlar ve kasa hareketleri gün sonu "
              "raporunda toplanır; PDF, Excel ve CSV olarak dışa aktarılır.",
              "Arrivals, departures, in-house guests and cash movements are "
              "collected in the night audit report and exported as PDF, Excel or "
              "CSV.")),
        ],
        "dipnot": ("Ayrı bir Finans ekranı henüz yok: tahsilat ve folyo işlemleri Ön "
                   "Büro'dan, mali özetler Raporlar ekranından yapılır.",
                   "There is no separate Finance screen yet: payments and folio work "
                   "happen on the Front Desk, financial summaries under Reports."),
    },

    # ── 12. KVKK / kimlik ────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🔐  Kimlik Verisi Şifreli Durur, Aranabilir Kalır",
                   "🔐  Identity Data Stays Encrypted — and Searchable"),
        "alt": ("Şifrelenmiş bir sütunda arama yapılamaz. Program bu ikilemi kör "
                "indeksle çözer: veri şifreli saklanır, arama özet değeri üzerinden yapılır.",
                "You cannot search an encrypted column. The program resolves that "
                "dilemma with a blind index: the data stays encrypted, the search "
                "runs over a keyed digest."),
        "kartlar": [
            (("🗄", "Alan Seviyesi Şifreleme", "Field-Level Encryption"),
             ("Kimlik ve pasaport numaraları veritabanına şifrelenerek yazılır. "
              "Doğrulama ham SQL sorgusuyla yapıldı: sütunda açık numara görünmüyor.",
              "National ID and passport numbers are written to the database "
              "encrypted. Verified with a raw SQL query: the column shows no "
              "readable number.")),
            (("🔎", "Kör İndeks", "Blind Index"),
             ("Arama, numaranın HMAC-SHA256 özetiyle yapılır: tam eşleşme çalışır, "
              "veri açığa çıkmaz. Kısmi arama bilinçli olarak desteklenmez.",
              "Search runs on an HMAC-SHA256 digest of the number: exact match "
              "works, the value is never exposed. Partial search is deliberately "
              "unsupported.")),
            (("👁", "Açık Görmek Ayrı Yetki", "Unmasking Is a Separate Right"),
             ("Numarayı maskesiz görmek ayrı bir izindir; kaydı görüntüleyebilen "
              "her kullanıcı bu izne sahip değildir ve görüntüleme denetime yazılır.",
              "Seeing the unmasked number is a separate permission; not every user "
              "who can view the record holds it, and each unmasking is written to "
              "the audit log.")),
            (("📇", "Misafir CRM", "Guest CRM"),
             ("Profil, konaklama geçmişi, tercihler, VIP ve kara liste, KVKK izin "
              "kayıtları. Özet şu an yalnız asıl misafiri sayar; refakatçi geceleri "
              "hariçtir.",
              "Profile, stay history, preferences, VIP and blacklist, consent "
              "records. The summary currently counts only the primary guest; "
              "companion nights are excluded.")),
        ],
        "dipnot": ("Loglarda API anahtarı, e-posta, telefon, kimlik ve kart numarası "
                   "otomatik maskelenir; maskeleme ayrı bir testle korunur.",
                   "API keys, e-mail addresses, phone, ID and card numbers are masked "
                   "automatically in logs; the masking is protected by its own test."),
    },

    # ── 13. Ekran: Misafir Kartı ve KVKK ─────────────────────────────────
    {
        "tip": "ekran",
        "baslik": ("👥  Ekranda: Misafir Kartı ve KVKK",
                   "👥  On Screen: Guest Profile and Privacy"),
        "alt": ("Kimlik numarası şifreli saklanır ve maskeli gösterilir; 'Göster' "
                "düğmesi her açık görüntülemeyi kullanıcı adı ve zaman damgasıyla "
                "denetim günlüğüne yazar.",
                "The ID number is stored encrypted and displayed masked; the 'Show' "
                "button writes every reveal to the audit log with the user name and "
                "a timestamp."),
        "gorseller": [("05_guests.png", None)],
    },

    # ── 14. Ölçüm ve raporlar ────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("📈  Ölçüm — Otelcilik Göstergeleri",
                   "📈  Measurement — Hotel Metrics"),
        "alt": ("Doluluk, ADR, RevPAR ve ALOS aynı sorgu katmanından üretilir: "
                "ekranda görünen sayı ile rapordaki sayı aynı koddan gelir.",
                "Occupancy, ADR, RevPAR and ALOS come from one query layer: the "
                "number on screen and the number in the report are produced by the "
                "same code."),
        "kartlar": [
            (("📊", "KPI ve Uyarılar", "KPIs and Alerts"),
             ("Doluluk, ADR, RevPAR, ALOS, iptal ve no-show oranı. Kritik uyarılar "
              "panelin üstünde toplanır: geciken çıkış, açık arıza, kritik stok.",
              "Occupancy, ADR, RevPAR, ALOS, cancellation and no-show rates. "
              "Critical alerts sit at the top: overdue departures, open faults, "
              "low stock.")),
            (("🔮", "14 Günlük Doluluk", "14-Day Occupancy"),
             ("Onaylı rezervasyonlardan önümüzdeki iki haftanın doluluk eğrisi "
              "çizilir. Bu bir talep tahmini modeli değil, bilinen rezervasyonların "
              "izdüşümüdür.",
              "The next two weeks are drawn from confirmed reservations. This is "
              "not a demand-forecasting model but a projection of bookings already "
              "on the books.")),
            (("🧷", "Kanal ve Oda Tipi", "Channel and Room Type"),
             ("Gelir kanal bazında (doğrudan, acente, çevrim içi) ve oda tipi "
              "bazında ayrıştırılır; hangi kanalın ne getirdiği tek tabloda görünür.",
              "Revenue is split by channel (direct, agency, online) and by room "
              "type; what each channel actually brings in is visible in a single "
              "table.")),
            (("📤", "PDF · Excel · CSV", "PDF · Excel · CSV"),
             ("Her rapor üç biçimde dışa aktarılır. CSV, Excel'in Türkçe "
              "karakterleri doğru açması için UTF-8 BOM ile yazılır.",
              "Every report exports in three formats. CSV is written with a UTF-8 "
              "BOM so that Excel opens accented characters correctly.")),
        ],
        "dipnot": ("106 raporlama testi çıktı biçimlerini ve boş veri durumunu da "
                   "kapsar: veri yokken rapor çökmez, boş tablo üretir.",
                   "106 reporting tests also cover export formats and the empty-data "
                   "case: with no data the report does not crash, it renders empty."),
    },

    # ── 15. Ekran: Raporlar ──────────────────────────────────────────────
    {
        "tip": "ekran",
        "baslik": ("📈  Ekranda: Raporlar", "📈  On Screen: Reports"),
        "alt": ("Doluluk, kanal ve oda tipi geliri, gün sonu kapanışı, kat hizmetleri "
                "ve KPI özeti; her rapor PDF, Excel ve CSV olarak dışa aktarılır.",
                "Occupancy, channel and room-type revenue, night audit, housekeeping "
                "and a KPI summary; every report can be exported to PDF, Excel and CSV."),
        "gorseller": [("10_reports.png", None)],
    },

    # ── 16. Yapay zekâ sınırı ────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🤖  Yapay Zekâ Öneri Üretir, Veriyi Değiştirmez",
                   "🤖  The AI Advises — It Never Writes"),
        "alt": ("Yapay zekâ katmanı veritabanı ile model arasındaki güvenlik "
                "sınırıdır: ne kişisel veri dışarı çıkar, ne de model kendiliğinden kayıt değiştirir.",
                "The AI layer is the security boundary between the database and the "
                "model: no personal data leaves, and the model never changes a "
                "record on its own."),
        "kartlar": [
            (("🛑", "İstem Denetimi", "Prompt Screening"),
             ("Modele gidecek metindeki e-posta, telefon ve uzun numara dizileri "
              "maskelenir; maskelemeden sonra kişisel veri kaldıysa çağrı hiç yapılmaz.",
              "E-mail addresses, phone numbers and long digit strings are masked "
              "before the call; if any personal data remains after masking, the "
              "request is never sent.")),
            (("🏠", "Yerel Seçenek", "Local Option"),
             ("LM Studio ile model kendi bilgisayarınızda çalışır: bulut yok, API "
              "ücreti yok, veri makineden çıkmaz. Bulut sağlayıcıları isteğe bağlıdır.",
              "With LM Studio the model runs on your own machine: no cloud, no API "
              "bill, no data leaving the building. Cloud providers are optional.")),
            (("💡", "Yalnız Öneri", "Advice Only"),
             ("Fiyat önerisi, misafir yorumu duygu analizi ve mesaj taslağı "
              "üretilir. Hiçbiri kendiliğinden uygulanmaz; karar kullanıcıdadır.",
              "Rate suggestions, guest-review sentiment analysis and message drafts "
              "are produced. None is applied automatically; the decision stays with "
              "the user.")),
            (("🏷", "Rozet ve Maliyet", "Badge and Cost"),
             ("Her çıktı \"AI tarafından oluşturuldu\" rozetiyle işaretlenir; jeton "
              "sayısı, süre ve tahmini maliyet gösterilir — yerel modelde ücretsiz.",
              "Every output is marked with an \"AI generated\" badge; token count, "
              "duration and estimated cost are shown — free on a local model.")),
        ],
        "dipnot": ("Birincil sağlayıcı düşerse yalnız GEÇİCİ hatalarda yedeğe geçilir. "
                   "Geçersiz API anahtarı kalıcı hatadır: gizlenmez, kullanıcıya söylenir.",
                   "If the primary provider fails, fallback happens only on TRANSIENT "
                   "errors. An invalid API key is permanent: it is reported, not hidden."),
    },

    # ── 17. Ekran: Yapay Zekâ Merkezi ────────────────────────────────────
    {
        "tip": "ekran",
        "baslik": ("🤖  Ekranda: Yapay Zekâ Merkezi", "🤖  On Screen: The AI Center"),
        "alt": ("Hazır görevler ve serbest soru; sağlayıcı ile model tek ekrandan "
                "seçilir. Yanıtlar 'AI tarafından oluşturuldu' rozetiyle işaretlenir.",
                "Ready-made tasks and free-form questions; provider and model are "
                "chosen on the same screen. Every answer carries an 'AI-generated' "
                "badge."),
        "gorseller": [("11_ai_center.png", None)],
    },

    # ── 18. AI Geliştirme Merkezi ────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🛠  AI Geliştirme Merkezi — Kısıtlanmış Ortam",
                   "🛠  AI Development Centre — A Restricted Sandbox"),
        "alt": ("Yapay zekâya kodlama görevi verilebilir; ama yapabilecekleri komut "
                "politikası, sandbox ve kalite zinciriyle sınırlanmıştır.",
                "The AI can be given coding tasks; what it may do is bounded by a "
                "command policy, a sandbox and a quality chain."),
        "kartlar": [
            (("📜", "Komut Politikası", "Command Policy"),
             ("İki katman: (1) izin listesi hangi komutun çalışacağına, (2) hedef "
              "denetimi hangi dosyaya dokunulabileceğine karar verir. `.env`, "
              "anahtar dosyaları ve misafir veritabanı — okuyucu hangisi olursa "
              "olsun — kapalıdır.",
              "Two layers: (1) an allow-list decides which command may run, "
              "(2) a target check decides which file it may touch. `.env`, key "
              "files and the guest database are closed to every reader, whichever "
              "one is used.")),
            (("📦", "Sandbox ve Diff", "Sandbox and Diff"),
             ("Yazma yalnız sandbox kökünde yapılır; her değişiklik önce diff "
              "olarak gösterilir, onaysız yazılmaz ve kayıp güncelleme koruması vardır.",
              "Writes happen only under the sandbox root; every change is shown as "
              "a diff first, nothing is written without approval, and lost-update "
              "protection is in place.")),
            (("🌿", "Git Koruması", "Git Guard"),
             ("Görev ayrı dalda çalışır, başlamadan kontrol noktası alınır, "
              "başarısızlıkta otomatik geri alınır. Ana dala birleştirme ayrı onay ister.",
              "The task runs on its own branch, a checkpoint is taken before it "
              "starts, and failure rolls back automatically. Merging to main needs "
              "separate approval.")),
            (("✅", "Kalite Zinciri", "Quality Chain"),
             ("Biçim → lint → tip → test → güvenlik. Testler geçmeden değişiklik "
              "işlenmez ve bu kural kullanıcı tarafından gevşetilemez.",
              "Format → lint → types → tests → security. No change is committed "
              "until the tests pass, and the user cannot relax that rule.")),
        ],
        "dipnot": ("__DEVTEST__ test bu modülü kapsar. Komutun fiilen çalıştırıldığı "
                   "katmanın kapsamı düşüktür ve bu boşluk test raporunda açıkça "
                   "yazılıdır — kapsam sayıları docs/TEST_REPORT.md'de ölçümüyle durur.",
                   "__DEVTEST__ tests cover this module. The layer that actually executes "
                   "the command has low coverage; that gap is stated openly in the test "
                   "report, with its measurement."),
    },

    # ── 19. Güvenlik ─────────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🛡  Güvenlik — Ölçülmüş, Varsayılmamış",
                   "🛡  Security — Measured, Not Assumed"),
        "alt": ("Güvenlik incelemesi bir iddia listesi değil: her madde bir testin "
                "ya da çalıştırılmış bir aracın çıktısına dayanır.",
                "The security review is not a list of claims: each item rests on a "
                "test or on the output of a tool that was actually run."),
        "kartlar": [
            (("🔑", "Kimlik Doğrulama", "Authentication"),
             ("Parolalar Argon2id ile hash'lenir; kaba kuvvet kilidi, oturum zaman "
              "aşımı ve kullanıcı sayımını engelleyen tek tip hata mesajı vardır.",
              "Passwords are hashed with Argon2id; brute-force lockout, session "
              "timeout and a uniform error message that prevents user enumeration "
              "are in place.")),
            (("🧩", "İki Katmanlı Yetki", "Authorization Twice"),
             ("Yetki hem arayüzde (menü, düğme) hem servis katmanında denetlenir. "
              "Ekranı gizlemek yetmez; yetkisiz servis çağrısı da reddedilir.",
              "Permissions are checked both in the UI (menu, buttons) and in the "
              "service layer. Hiding a screen is not enough; the unauthorized "
              "service call is rejected too.")),
            (("🗝", "Sır Yönetimi", "Secret Handling"),
             ("API anahtarları Windows Credential Manager'da tutulur; veritabanına, "
              "kaynak koda ve hata mesajlarına yazılmaz. Depoda sır olmadığı tarandı.",
              "API keys live in Windows Credential Manager; they are never written "
              "to the database, the source or error messages. The repository was "
              "scanned for secrets.")),
            (("📚", "Denetim Günlüğü", "Audit Log"),
             ("Kritik işlemler yalnız eklenebilir bir denetim günlüğüne yazılır: "
              "kim, ne zaman, hangi kayıtta ne yaptı.",
              "Critical operations are written to an append-only audit log: who did "
              "what, when, and to which record.")),
        ],
        "dipnot": ("bandit: __BANDIT_YUKSEK__ yüksek, __BANDIT_ORTA__ orta, "
                   "__BANDIT_DUSUK__ düşük önemde bulgu — düşük bulguların her biri "
                   "gerekçelendirilmiştir. Sayılar sunum üretilirken ölçülür.",
                   "bandit: __BANDIT_YUKSEK__ high, __BANDIT_ORTA__ medium and "
                   "__BANDIT_DUSUK__ low-severity findings — each low finding carries a "
                   "written rationale. The numbers are measured as the deck is built."),
    },

    # ── 20. Bulunan açıklar ──────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🔍  İncelemeler Gerçek Açıklar Buldu",
                   "🔍  The Reviews Found Real Holes"),
        "alt": ("Bir güvenlik incelemesinin değeri bulduklarındadır. Aşağıdakilerin "
                "hepsi düzeltildi; her biri için bir gerileme testi vardır ve "
                "düzeltmeler docs/SECURITY_REVIEW.md 8. bölümde anlatılır.",
                "The value of a security review is in what it finds. Everything below "
                "was fixed, each has a regression test, and the fixes are described in "
                "docs/SECURITY_REVIEW.md section 8."),
        "kartlar": [
            (("🚪", "Yetki Boşluğu", "Missing Permission Check"),
             ("Odayı satışa kapatma işlemi yetki denetiminden geçmiyordu: bakım "
              "yetkisi olmayan bir kullanıcı odayı kapatabiliyordu. Kapatıldı; "
              "`tests/application/test_operations_services.py` koruyor.",
              "Taking a room out of sale skipped the permission check: a user "
              "without maintenance rights could block a room. Closed; guarded by "
              "`tests/application/test_operations_services.py`.")),
            (("💬", "Sahiplik Kontrolü Yoktu", "No Ownership Check"),
             ("Yapay zekâ sohbet geçmişinde sahiplik denetlenmiyordu; bir kullanıcı "
              "başkasının sohbetini okuyabilirdi. Kapatıldı; "
              "`tests/application/test_ai_service.py` koruyor.",
              "AI chat history did not verify ownership; one user could read "
              "another's conversation. Closed; guarded by "
              "`tests/application/test_ai_service.py`.")),
            (("🧨", "Politika Atlatmaları", "Policy Bypasses"),
             ("Komut politikasında atlatma yolları bulundu ve kapatıldı. Sonuncusu "
              "bağımsız bir yayın öncesi incelemede çıktı: izin listesindeki bir "
              "okuyucu (`head`, `findstr`) `.env` dosyasını onaysız okuyabiliyordu. "
              "Artık karar komutun adına değil dokunduğu dosyaya bakıyor.",
              "Bypasses were found in the command policy and closed. The last one "
              "came out of an independent pre-release review: an allow-listed reader "
              "(`head`, `findstr`) could read `.env` with no approval. The decision "
              "now looks at the file touched, not at the command's name.")),
            (("💰", "Mali Tutarsızlık", "Financial Inconsistency"),
             ("Folyo ve tahsilat aritmetiğindeki hatalar düzeltildi; yedeğe "
              "geçişte anahtar hatasının gizlenmesi de giderildi. Para yolunda "
              "`float` kullanılmaz, `Decimal` kullanılır — bu da test edilir.",
              "Arithmetic errors in folio and payment handling were corrected, and a "
              "key error being swallowed during provider fallback was fixed. Money "
              "never travels as `float`, only as `Decimal` — and that is tested.")),
        ],
        "dipnot": ("Bilerek işlem (commit) numarası verilmiyor: bir düzeltmenin "
                   "kanıtı onu koruyan testtir. Aynı belge düzeltilmeyenleri de "
                   "listeler — numara üreteçleri eşzamanlı güvenli değildir.",
                   "No commit hashes are cited on purpose: the evidence for a fix is "
                   "the test that guards it. The same document lists what was not "
                   "fixed — the number generators are not concurrency-safe."),
    },

    # ── 21. Test raporu ──────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🧪  Sayılar Hatırlanmadı, Çalıştırıldı",
                   "🧪  The Numbers Were Run, Not Remembered"),
        "alt": ("__TEST__ test geçiyor ve hiçbiri ağa çıkmıyor — ağa çıkılmadığı "
                "da ayrı bir testle korunuyor. Sayı bu sunum üretilirken ölçüldü.",
                "__TEST__ tests pass and none of them touches the network — and that "
                "fact is itself protected by a test. The figure was measured while "
                "this deck was generated."),
        "kartlar": [
            (("📦", "Katman Katman", "Layer by Layer"),
             ("__KATMAN__.",
              "__KATMAN__.")),
            (("🎯", "13 Kritik Senaryo", "13 Critical Scenarios"),
             ("Çakışan rezervasyon, iptal, no-show, erken giriş, kapalı oda, hatalı "
              "ödeme, yetkisiz erişim, model çökmesi, kota, geçersiz anahtar, bozuk "
              "JSON, boş rapor — hepsi kapsanıyor.",
              "Overlapping booking, cancellation, no-show, early arrival, blocked "
              "room, wrong payment, unauthorized access, dead model, quota, invalid "
              "key, malformed JSON, empty report — all covered.")),
            (("📉", "Boşluklar Yazılı", "Gaps Written Down"),
             ("Açılış yolu (giriş ekranı, ana pencere) ve komut satırı kapsam dışı "
              "kalır; yedekleme düşük kapsamlıdır. Bunlar gizlenmedi — "
              "docs/TEST_REPORT.md içinde ölçümleriyle listelendi ve "
              "gerekçelendirildi.",
              "The startup path (login screen, main window) and the CLI stay "
              "uncovered; backup has low coverage. None of this is hidden — it is "
              "listed with its measurements and explained in docs/TEST_REPORT.md.")),
            (("⚙️", "Her Push'ta CI", "CI on Every Push"),
             ("GitHub Actions (.github/workflows/ci.yml): biçim, lint, tip, test, "
              "güvenlik taraması, göç tutarlılığı ve \"depoya hassas dosya girdi mi\" "
              "denetimi. `ruff check app tests` şu an __RUFF__ bulgu veriyor.",
              "GitHub Actions (.github/workflows/ci.yml): format, lint, types, tests, "
              "security scanning, migration consistency and a \"did a sensitive file "
              "enter the repo\" gate. `ruff check app tests` currently reports "
              "__RUFF__ findings.")),
        ],
        "dipnot": ("Rapor bir anlık görüntüdür: __TARIH__, sürüm __SURUM__, "
                   "Windows 11, Python 3.11. Sayılar üretim anında ölçülür; "
                   "kod değişince değişirler.",
                   "The report is a snapshot: __TARIH__, version __SURUM__, "
                   "Windows 11, Python 3.11. The numbers are measured at generation "
                   "time and change when the code does."),
    },

    # ── 22. Mimari ───────────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🏗  Mimari — Bağımlılık Tek Yöne Akar",
                   "🏗  Architecture — Dependencies Flow One Way"),
        "alt": ("ui → application → domain ← infrastructure. İş kuralları ne "
                "veritabanını ne arayüzü tanır; 99 alan testinin saniyeler içinde "
                "koşmasının sebebi budur.",
                "ui → application → domain ← infrastructure. The business rules know "
                "neither the database nor the UI, which is why the 99 domain tests "
                "run in seconds."),
        "kartlar": [
            (("🧱", "Çerçeveden Bağımsız Alan", "Framework-Free Domain"),
             ("Alan katmanı veritabanı ve arayüz kitaplıklarını içe aktarmaz. Kural "
              "belgede yazmakla kalmaz, içe aktarım denetimiyle sınanır.",
              "The domain layer imports neither the ORM nor the UI toolkit. The rule "
              "is not only documented but enforced by an import test.")),
            (("🕐", "Zaman Dilimi Bilinçli", "Timezone-Aware"),
             ("Özel bir sütun tipi tüm zaman damgalarını UTC olarak saklar. Saat "
              "farkı ve yaz saati hataları veri katmanında değil, tek noktada çözülür.",
              "A custom column type stores every timestamp in UTC. Offset and "
              "daylight-saving bugs are solved in one place instead of scattered "
              "through the data layer.")),
            (("🔁", "Göç Edilebilir Şema", "Migratable Schema"),
             ("__TABLO__ tablo Alembic ile sürümlenir ve CI her değişiklikte şema ile "
              "göçlerin tutarlılığını doğrular. Durumlar sayı değil, metin olarak saklanır.",
              "60 tables are versioned with Alembic and CI verifies schema/migration "
              "consistency on every change. Enumerations are stored as text, not "
              "as integers.")),
            (("🧩", "Sayfa Kayıt Defteri", "Page Registry"),
             ("Yeni ekran eklemek için ana pencere değiştirilmez, kayıt defterine "
              "bir satır eklenir. Tamamlanmamış ekran da aynı defterden gösterilir.",
              "Adding a screen does not touch the main window: one line is added to "
              "the registry. Unfinished screens are declared in the same place.")),
        ],
        "dipnot": ("SQLite varsayılan, PostgreSQL destekleniyor. Türkçe büyük/küçük "
                   "harf araması SQLite'ta sınırlıdır; \"İ/ı\" için PostgreSQL gerekir.",
                   "SQLite by default, PostgreSQL supported. Case-insensitive Turkish "
                   "search is limited on SQLite; dotted/dotless \"I\" needs PostgreSQL."),
    },

    # ── 23. Kurulum ve işletme ───────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("⚙️  Kurulum, Yedek ve Paketleme",
                   "⚙️  Setup, Backup and Packaging"),
        "alt": ("Tek komutla kurulur, tek komutla yedeklenir, tek dosyaya "
                "paketlenir. Veritabanı ve loglar uygulamanın yanında kalır — taşınabilir.",
                "One command to install, one to back up, one to package. The database "
                "and logs live next to the application, so the installation stays "
                "portable."),
        "kartlar": [
            (("🚀", "Tek Komut Kurulum", "One-Command Setup"),
             ("Kurulum betiği Python sürümünü doğrular, sanal ortamı kurar, .env "
              "dosyasını rastgele oturum anahtarıyla üretir, göçleri uygular ve "
              "yönetici hesabını açar.",
              "The setup script validates the Python version, creates the virtual "
              "environment, writes .env with a random session key, applies "
              "migrations and creates the admin account.")),
            (("💾", "Tutarlı Yedek", "Consistent Backup"),
             ("SQLite yedeği VACUUM INTO ile alınır; WAL kipinde dosyayı kopyalamak "
              "tutarsız yedek üretir. Geri yüklemeden önce mevcut veritabanı saklanır.",
              "The SQLite backup uses VACUUM INTO; copying the file in WAL mode "
              "produces an inconsistent backup. The current database is kept aside "
              "before any restore.")),
            (("📦", "Platformlar ve İndirme", "Platforms and Download"),
             ("Windows 10/11 (x64) + macOS (Apple Silicon) — İndirme: __DEPO__/"
              "releases (v__SURUM__). macOS paketi notarize edilmedi; ilk açılış: "
              "sağ tık → Aç.",
              "Windows 10/11 (x64) & macOS (Apple Silicon) — Download: GitHub "
              "Releases (v__SURUM__). The macOS package is not notarized; first "
              "launch: right-click → Open.")),
            (("🧰", "Bakım Komutları", "Maintenance Commands"),
             ("bootstrap · seed-demo · backup · restore · check-ai · doctor. "
              "check-ai sağlayıcı bağlantısını, doctor kurulumun sağlığını sınar.",
              "bootstrap · seed-demo · backup · restore · check-ai · doctor. "
              "check-ai tests the provider connection, doctor the health of the "
              "installation.")),
        ],
        "dipnot": ("Demo kurulumda beş rol hesabı hazır gelir (müdür, ön büro, kat, "
                   "teknik, muhasebe). Gerçek kurulumda demo veri ve hesaplar silinmelidir.",
                   "A demo installation ships five role accounts (manager, front desk, "
                   "housekeeping, maintenance, accounting). In a real deployment demo "
                   "data and accounts must be removed."),
    },

    # ── 24. Yapılmayanlar ────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("📋  Yapılmayanlar — Aynı Netlikte",
                   "📋  What Is Not Done — Stated Just as Clearly"),
        "alt": ("Bir özelliğin yol haritasında olması, bugün üretimde "
                "kullanılmaması gerektiği anlamına gelir. Program eksiğini kullanıcıdan saklamaz.",
                "If a feature is on the roadmap, it means it should not be used in "
                "production today. The program does not hide its gaps from the user."),
        "kartlar": [
            (("🧾", "e-Fatura ve Bildirim", "e-Invoice and Notification"),
             ("YAPILMADI. Fatura tablosunda alanlar, ayar dosyasında anahtarlar "
              "duruyor; kurum entegrasyonu, XML üretimi ve bildirim akışı yok.",
              "NOT DONE. The invoice table has the fields and the settings file the "
              "keys; the authority integration, XML generation and notification flow "
              "do not exist.")),
            (("🖥", "Finans · Stok · Personel", "Finance · Stock · Staff"),
             ("İş mantığı ve raporlar hazır, ayrı ekran yok. Ekran açıldığında ne "
              "planlandığı ve aynı işin şu anda nasıl yapılacağı yazılı olarak gösterilir.",
              "Business logic and reports exist, dedicated screens do not. Opening "
              "one shows what is planned and how to do the same job today.")),
            (("🔢", "Numara Üreteci", "Number Generator"),
             ("Rezervasyon ve folyo numaraları en büyük değere bir eklenerek "
              "üretilir; eşzamanlı kullanımda ikinci işlem reddedilir ve yeniden "
              "denenmelidir. Çözüm v0.2'de.",
              "Reservation and folio numbers are produced with max()+1; under "
              "concurrent use the second write is rejected and must be retried. The "
              "fix lands in v0.2.")),
            (("☁️", "Bulut Sağlayıcıları", "Cloud Providers"),
             ("NVIDIA ve Anthropic sağlayıcı kodu hazır ama gerçek API çağrısıyla "
              "denenmedi — anahtar gerektirir. LM Studio gerçek istekle doğrulandı.",
              "The NVIDIA and Anthropic adapters are written but never exercised "
              "against the real API — that needs a key. LM Studio was verified with "
              "a real request.")),
        ],
        "dipnot": ("Belge soru-cevabı (RAG) ve görsel belge analizi de tamamlanmadı: "
                   "veri modeli hazır, iş akışı bağlanmadı.",
                   "Document question-answering (RAG) and visual document analysis are "
                   "also unfinished: the data model exists, the workflow is not wired."),
    },

    # ── 25. Yol haritası ─────────────────────────────────────────────────
    {
        "tip": "kart",
        "baslik": ("🗺  Yol Haritası", "🗺  Roadmap"),
        "alt": ("Sıra rastgele değil: önce günlük işi tamamlayan ekranlar, sonra "
                "yapay zekâ derinliği, en sonda mevzuat entegrasyonları.",
                "The order is deliberate: first the screens that complete daily "
                "operations, then depth in AI, and regulatory integrations last."),
        "kartlar": [
            (("🧭", "v0.2 — Operasyon", "v0.2 — Operations"),
             ("Finans, Stok ve oda tipi/fiyat planı ekranları; personel ve vardiya "
              "yönetimi; numara üretiminde eşzamanlılık düzeltmesi; grup rezervasyonu "
              "arayüzü.",
              "Finance, stock and room-type/rate-plan screens; staff and shift "
              "management; the concurrency fix for number generation; a group "
              "booking interface.")),
            (("🧠", "v0.3 — Yapay Zekâ", "v0.3 — Artificial Intelligence"),
             ("Belge indeksleme akışının tamamlanması, görsel belge analizi (KVKK "
              "değerlendirmesiyle), talep tahmini ve dinamik fiyat önerisi.",
              "Completing the document indexing flow, visual document analysis (with "
              "a data-protection assessment), demand forecasting and dynamic rate "
              "suggestions.")),
            (("🔌", "v0.4 — Entegrasyon", "v0.4 — Integrations"),
             ("e-Fatura entegratör bağlantısı, kimlik bildirimi, kanal yöneticisi "
              "taslağı ve e-posta/SMS bildirim altyapısı.",
              "The e-invoice integrator connection, guest notification, a channel "
              "manager draft and the e-mail/SMS notification layer.")),
            (("🏁", "v1.0 — Üretim", "v1.0 — Production"),
             ("PostgreSQL üretim kurulumu, çok tesisli test, bağımsız güvenlik "
              "denetimi, 10.000+ rezervasyonla performans testi, kullanıcı kabul testleri.",
              "PostgreSQL production setup, multi-property testing, an independent "
              "security audit, performance testing with 10,000+ reservations and user "
              "acceptance tests.")),
        ],
        "dipnot": ("Katkıya en açık alanlar: mevzuat entegrasyonları, gerçek otel "
                   "işletmelerinden iş akışı geri bildirimi ve İngilizce arayüz çevirisi.",
                   "Most open to contribution: regulatory integrations, workflow "
                   "feedback from real hotels, and the English interface translation."),
    },

    # ── 26. Kapanış ──────────────────────────────────────────────────────
    {
        "tip": "kapak",
        "simge": True,
        "pul": False,
        "baslik": ("NE SATTIĞINIZI BİLEREK YÖNETİN",
                   "RUN THE HOUSE KNOWING WHAT YOU SOLD"),
        "alt": ("__TABLO__ tablo · __TEST__ test · __KAPSAM__ kapsam · şifreli "
                "kimlik verisi · yerel yapay zekâ seçeneği",
                "__TABLO__ tables · __TEST__ tests · __KAPSAM__ coverage · encrypted "
                "identity data · a local AI option"),
        "slogan": ("Eksikleri de yazılı olan bir program.",
                   "A program that writes down its gaps too."),
        "puller": [
            ("Kullanıcı kılavuzu", "User guide"),
            ("Test raporu", "Test report"),
            ("Güvenlik incelemesi", "Security review"),
            ("Yol haritası", "Roadmap"),
        ],
        "bilgi": ("__DEPO__ · MIT lisansı · Sürüm __SURUM__ · __AY__",
                  "__DEPO__ · MIT license · Version __SURUM__ · __AY__"),
    },
]


# ══════════════════════════════════════════════════════════════════════════
# ÇİZİM YARDIMCILARI
# ══════════════════════════════════════════════════════════════════════════

def _renk(hex6):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex6.upper())


def _inc(deger):
    from pptx.util import Inches
    return Inches(deger)


def _kutu_ekle(slayt, kutu, *, dolgu, cerceve=None, oval=False, yuvarlaklik=0.10):
    """Yuvarlatılmış dikdörtgen (ya da daire) ekler.

    Gölge açıkça KAPATILIR: python-pptx yeni şekle temanın gölgesini miras
    bırakır ve koyu zeminde bu, kartların altında bulanık bir leke bırakır.
    """
    from pptx.enum.shapes import MSO_SHAPE
    l, t, g, y = kutu
    bicim = MSO_SHAPE.OVAL if oval else MSO_SHAPE.ROUNDED_RECTANGLE
    sekil = slayt.shapes.add_shape(bicim, _inc(l), _inc(t), _inc(g), _inc(y))
    if not oval:
        try:
            sekil.adjustments[0] = yuvarlaklik
        except Exception:
            pass
    sekil.fill.solid()
    sekil.fill.fore_color.rgb = _renk(dolgu)
    if cerceve:
        from pptx.util import Pt
        sekil.line.color.rgb = _renk(cerceve)
        sekil.line.width = Pt(1)
    else:
        sekil.line.fill.background()
    try:
        sekil.shadow.inherit = False
    except Exception:
        pass
    sekil.text_frame.word_wrap = True
    return sekil


def _metin_ekle(slayt, kutu, parcalar, *, hiza="sol", dikey="orta", satir_araligi=None):
    """Metin kutusu ekler. `parcalar` = [(metin, punto, kalın, renk_hex)].

    Tek şekilde birden çok run kullanılabilmesi başlıklar için gerekli:
    ikon 28 punto normal, başlık 30 punto kalın — ikisi aynı satırda ama
    aynı biçimde değil.
    """
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt
    l, t, g, y = kutu
    kutucuk = slayt.shapes.add_textbox(_inc(l), _inc(t), _inc(g), _inc(y))
    tf = kutucuk.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"ust": MSO_ANCHOR.TOP, "orta": MSO_ANCHOR.MIDDLE}[dikey]
    p = tf.paragraphs[0]
    p.alignment = {"sol": PP_ALIGN.LEFT, "orta": PP_ALIGN.CENTER}[hiza]
    if satir_araligi:
        p.line_spacing = satir_araligi
    for metin, punto, kalin, renk in parcalar:
        run = p.add_run()
        run.text = metin
        run.font.size = Pt(punto)
        run.font.bold = kalin
        run.font.name = YAZI_TIPI
        run.font.color.rgb = _renk(renk)
    return kutucuk


def _zemin(slayt, palet):
    """Slayt arka planını boyar.

    Arka plan bir şekil değildir; slaydın kendi tanımında durur. Atlanırsa
    düzenin varsayılan BEYAZ zemini kalır ve koyu kartlar beyaz sayfada
    yüzer — yani sunum iki tasarımlı görünür.
    """
    dolgu = slayt.background.fill
    dolgu.solid()
    dolgu.fore_color.rgb = _renk(palet["zemin"])


def _baslik_bandi(slayt, palet, baslik):
    """Başlığı ikon + metin olarak iki run hâlinde yazar."""
    if "  " in baslik:
        ikon, kalan = baslik.split("  ", 1)
        parcalar = [(ikon + "  ", P_BASLIK_IKON, False, palet["yazi"]),
                    (kalan, P_BASLIK, True, palet["yazi"])]
    else:
        parcalar = [(baslik, P_BASLIK, True, palet["yazi"])]
    _metin_ekle(slayt, BASLIK_KUTU, parcalar)


# ══════════════════════════════════════════════════════════════════════════
# SLAYT KURUCULARI
# ══════════════════════════════════════════════════════════════════════════

def _kapak_slaydi(sunum, palet, veri, dil, surum):
    slayt = sunum.slides.add_slide(sunum.slide_layouts[6])
    _zemin(slayt, palet)

    # Kapak ile kapanış aynı kuruluşu paylaşır: kapanışta sürüm pulu yoktur
    # (alt bilgide zaten yazılı) ve başlık, pulun boşalttığı yere kayar.
    simge_var = veri.get("simge") and os.path.exists(SIMGE)
    if simge_var:
        boy = 1.20
        slayt.shapes.add_picture(SIMGE, _inc(ORTA - boy / 2), _inc(0.75),
                                 _inc(boy), _inc(boy))
    if veri.get("pul"):
        _kutu_ekle(slayt, (ORTA - 0.85, 2.12, 1.70, 0.42),
                   dolgu=palet["kart"], cerceve=palet["cerceve"], yuvarlaklik=0.50)
        _metin_ekle(slayt, (ORTA - 0.85, 2.12, 1.70, 0.42),
                    [("v" + surum, P_PUL, True, palet["marka"])], hiza="orta")
        y_baslik = 2.72
    else:
        y_baslik = 2.40 if simge_var else 2.20

    _metin_ekle(slayt, (0.60, y_baslik, 12.13, 0.95),
                [(_M(veri["baslik"], dil), P_KAPAK_BASLIK, True, palet["yazi"])],
                hiza="orta")
    _metin_ekle(slayt, (1.00, y_baslik + 1.05, 11.33, 0.70),
                [(_M(veri["alt"], dil), P_KAPAK_ALT, False, palet["marka"])],
                hiza="orta")
    _metin_ekle(slayt, (1.00, y_baslik + 1.85, 11.33, 0.50),
                [(_M(veri["slogan"], dil), P_KAPAK_SLOGAN, False, palet["altin"])],
                hiza="orta")

    for indis, pul_metni in enumerate(veri["puller"]):
        kutu = (PUL_X[indis], 5.30, PUL_OLCU[0], PUL_OLCU[1])
        _kutu_ekle(slayt, kutu, dolgu=palet["kart"], cerceve=palet["cerceve"],
                   yuvarlaklik=0.50)
        _metin_ekle(slayt, kutu, [(_M(pul_metni, dil), P_PUL, False, palet["yazi"])],
                    hiza="orta")

    _metin_ekle(slayt, (1.00, 6.62, 11.33, 0.40),
                [(_M(veri["bilgi"], dil), P_DIPNOT, False, palet["soluk"])],
                hiza="orta")
    return slayt


def _sayi_slaydi(sunum, palet, veri, dil):
    slayt = sunum.slides.add_slide(sunum.slide_layouts[6])
    _zemin(slayt, palet)
    _baslik_bandi(slayt, palet, _M(veri["baslik"], dil))

    for indis, (rakam, rol, aciklama) in enumerate(veri["kutular"]):
        sutun, satir = indis % 4, indis // 4
        l, t = SAYI_X[sutun], SAYI_T[satir]
        _kutu_ekle(slayt, (l, t, SAYI_G, SAYI_Y), dolgu=palet["kart"],
                   cerceve=palet["cerceve"])
        _metin_ekle(slayt, (l + RAKAM_OFS[0], t + RAKAM_OFS[1],
                            RAKAM_OLCU[0], RAKAM_OLCU[1]),
                    [(_M(rakam, dil) if isinstance(rakam, tuple) else rakam,
                      P_RAKAM, True, palet[rol])], hiza="orta")
        _metin_ekle(slayt, (l + ACIKLAMA_OFS[0], t + ACIKLAMA_OFS[1],
                            ACIKLAMA_OLCU[0], ACIKLAMA_OLCU[1]),
                    [(_M(aciklama, dil), P_ACIKLAMA, False, palet["soluk"])],
                    hiza="orta", satir_araligi=0.95)

    _kutu_ekle(slayt, NOT_KUTU, dolgu=palet["kart"], cerceve=palet["cerceve"])
    _metin_ekle(slayt, (NOT_KUTU[0] + 0.30, NOT_KUTU[1] + 0.15,
                        NOT_KUTU[2] - 0.60, NOT_KUTU[3] - 0.30),
                [(_M(veri["not"], dil), P_ALT, False, palet["soluk"])])
    return slayt


def _kart_slaydi(sunum, palet, veri, dil):
    slayt = sunum.slides.add_slide(sunum.slide_layouts[6])
    _zemin(slayt, palet)
    _baslik_bandi(slayt, palet, _M(veri["baslik"], dil))
    _metin_ekle(slayt, ALT_KUTU,
                [(_M(veri["alt"], dil), P_ALT, False, palet["soluk"])],
                dikey="ust", satir_araligi=1.0)

    for indis, ((ikon, ad_tr, ad_en), govde) in enumerate(veri["kartlar"]):
        sutun, satir = indis % 2, indis // 2
        l, t = KART_X[sutun], KART_T[satir]
        _kutu_ekle(slayt, (l, t, KART_G, KART_Y), dolgu=palet["kart"],
                   cerceve=palet["cerceve"], yuvarlaklik=0.06)
        daire = (l + IKON_OFS[0], t + IKON_OFS[1], IKON_CAP, IKON_CAP)
        _kutu_ekle(slayt, daire, dolgu=palet[DAIRE_DONGUSU[indis % 4]], oval=True)
        _metin_ekle(slayt, daire, [(ikon, P_KART_IKON, False, palet["yazi"])],
                    hiza="orta")
        _metin_ekle(slayt, (l + AD_OFS[0], t + AD_OFS[1], AD_OLCU[0], AD_OLCU[1]),
                    [((ad_tr if dil == "tr" else ad_en), P_KART_AD, True, palet["yazi"])])
        _metin_ekle(slayt, (l + GOVDE_OFS[0], t + GOVDE_OFS[1],
                            GOVDE_OLCU[0], GOVDE_OLCU[1]),
                    [(_M(govde, dil), P_KART_GOVDE, False, palet["soluk"])],
                    dikey="ust", satir_araligi=1.0)

    if veri.get("dipnot"):
        _metin_ekle(slayt, DIPNOT_KUTU,
                    [(_M(veri["dipnot"], dil), P_DIPNOT, False, palet["soluk"])],
                    dikey="ust")
    return slayt


def _gorsel_yolu(ad):
    """Ekran görüntüsünün tam yolu; dosya yoksa üretimi durdurur.

    Eksik görselle devam etmek, sunumda boş bir çerçeve basmak demektir —
    bunun yerine ne yapılacağını söyleyen bir hata verilir.
    """
    yol = os.path.join(EKRAN_DIZIN, ad)
    if not os.path.exists(yol):
        raise FileNotFoundError(
            "%s bulunamadı. Önce ekran görüntülerini üretin:\n"
            "    .venv\\Scripts\\python.exe sunum\\ekran_yakala.py" % yol)
    return yol


def _resim_cercevesi(resim, palet):
    """Görüntüye ince çerçeve çizer, tema gölgesini kapatır.

    Koyu arayüz görüntüsü koyu zeminde çerçevesiz kaybolur; baskı sürümünün
    beyaz zemininde ise çerçevesiz görüntü sayfaya yapışık durur.
    """
    from pptx.util import Pt
    resim.line.color.rgb = _renk(palet["cerceve"])
    resim.line.width = Pt(1)
    try:
        resim.shadow.inherit = False
    except Exception:
        pass


def _ekran_slaydi(sunum, palet, veri, dil):
    """Ekran görüntüsü slaydı: 1 büyük ya da 2 yan yana görsel.

    Görüntüler 16:9'dur; tek görselde yükseklik sabitlenir ve python-pptx
    genişliği orandan türetir, ardından yatay ortalanır. Çift görselde kart
    ızgarasının sütunları kullanılır ve üstlerine birer etiket yazılır.
    """
    slayt = sunum.slides.add_slide(sunum.slide_layouts[6])
    _zemin(slayt, palet)
    _baslik_bandi(slayt, palet, _M(veri["baslik"], dil))
    _metin_ekle(slayt, ALT_KUTU,
                [(_M(veri["alt"], dil), P_ALT, False, palet["soluk"])],
                dikey="ust", satir_araligi=1.0)

    gorseller = veri["gorseller"]
    if len(gorseller) == 1:
        dosya, _etiket = gorseller[0]
        resim = slayt.shapes.add_picture(_gorsel_yolu(dosya), _inc(0),
                                         _inc(EKRAN_TEK_UST),
                                         height=_inc(EKRAN_TEK_YUK))
        resim.left = int((sunum.slide_width - resim.width) / 2)
        _resim_cercevesi(resim, palet)
    else:
        for indis, (dosya, etiket) in enumerate(gorseller):
            l = KART_X[indis]
            _metin_ekle(slayt, (l, EKRAN_CIFT_ETIKET_T, KART_G, 0.35),
                        [(_M(etiket, dil), P_KART_AD, True, palet["yazi"])])
            resim = slayt.shapes.add_picture(_gorsel_yolu(dosya), _inc(l),
                                             _inc(EKRAN_CIFT_T),
                                             width=_inc(KART_G))
            _resim_cercevesi(resim, palet)

    if veri.get("dipnot"):
        _metin_ekle(slayt, DIPNOT_KUTU,
                    [(_M(veri["dipnot"], dil), P_DIPNOT, False, palet["soluk"])],
                    dikey="ust")
    return slayt


#: Ölçülemeyen bir değer yerine basılacak metin. Slaytta rakam yerine
#: "ölçülmedi" görmek, yanlış bir rakam görmekten iyidir.
YOK = {"tr": "ölçülmedi", "en": "not measured"}


def _degiskenler(metin, dil, surum, olcum=None):
    """Sunum metnindeki yer tutucuları doldurur.

    Ölçüm yer tutucuları (`__TEST__`, `__TABLO__`, ...) `sunum/olcum.py`
    çıktısından gelir. Bir değer ölçülemediyse yerine "ölçülmedi" yazılır —
    hiçbir koşulda hatırlanmış bir sayı basılmaz.
    """
    olcum = olcum or {}

    def _s(anahtar, bicim="%s"):
        deger = olcum.get(anahtar)
        if deger is None:
            return YOK[dil]
        return bicim % deger

    kapsam = olcum.get("kapsam")
    if kapsam is None:
        kapsam_metni = YOK[dil]
    elif dil == "tr":
        kapsam_metni = ("%%%.1f" % kapsam).replace(".", ",")
    else:
        kapsam_metni = "%.1f%%" % kapsam

    katman = olcum.get("katman") or {}
    katman_tr = " · ".join(
        "%s %d" % (ad, katman[kod])
        for kod, ad in (
            ("alan", "alan"), ("guvenlik", "güvenlik"), ("altyapi", "altyapı"),
            ("uygulama", "uygulama"), ("yapayzeka", "yapay zekâ"),
            ("arayuz", "arayüz"), ("raporlama", "raporlama"),
            ("devmerkezi", "geliştirme merkezi"),
        )
        if kod in katman
    ) or YOK["tr"]
    katman_en = " · ".join(
        "%s %d" % (ad, katman[kod])
        for kod, ad in (
            ("alan", "Domain"), ("guvenlik", "security"), ("altyapi", "infrastructure"),
            ("uygulama", "application"), ("yapayzeka", "AI"), ("arayuz", "UI"),
            ("raporlama", "reporting"), ("devmerkezi", "development centre"),
        )
        if kod in katman
    ) or YOK["en"]

    bandit = olcum.get("bandit") or {}
    return (metin.replace("__SURUM__", surum)
                 .replace("__AY__", YAYIN_AYI[dil])
                 .replace("__TARIH__", OLCUM_TARIHI[dil])
                 .replace("__DEPO__", DEPO)
                 .replace("__TEST__", _s("test"))
                 .replace("__TABLO__", _s("tablo"))
                 .replace("__MODEL__", _s("model"))
                 .replace("__IZIN__", _s("izin"))
                 .replace("__ROL__", _s("rol"))
                 .replace("__SAGLAYICI__", _s("saglayici"))
                 .replace("__EKRAN__", _s("ekran"))
                 .replace("__KAPSAM__", kapsam_metni)
                 .replace("__KATMAN__", katman_tr if dil == "tr" else katman_en)
                 .replace("__DEVTEST__", str(katman.get("devmerkezi", YOK[dil])))
                 .replace("__BANDIT_YUKSEK__", str(bandit.get("HIGH", YOK[dil])))
                 .replace("__BANDIT_ORTA__", str(bandit.get("MEDIUM", YOK[dil])))
                 .replace("__BANDIT_DUSUK__", str(bandit.get("LOW", YOK[dil])))
                 .replace("__RUFF__", _s("ruff")))


def _icerik(surum):
    """Yer tutucuları doldurulmuş slayt verisi — iki dilli yapı korunur.

    Her dil kendi yer tutucusunu alır: `__AY__` Türkçede "Ağustos 2026",
    İngilizcede "August 2026" olur. Dil seçimi çizim anında `_M` ile yapılır.
    """
    olcum = olcumler()
    cozulmus = []
    for slayt in SLAYTLAR:
        kopya = dict(slayt)
        for anahtar in ("baslik", "alt", "slogan", "dipnot", "not", "bilgi"):
            if anahtar in kopya:
                tr, en = kopya[anahtar]
                kopya[anahtar] = (_degiskenler(tr, "tr", surum, olcum),
                                  _degiskenler(en, "en", surum, olcum))
        if kopya.get("tip") == "sayi" and "kutular" in kopya:
            # Büyük rakam da iki dilli çözülür: ondalık ayıracı ve "ölçülmedi"
            # metni dile göre değişir (%77,6 / 77.6%).
            kopya["kutular"] = [
                ((_degiskenler(buyuk, "tr", surum, olcum),
                  _degiskenler(buyuk, "en", surum, olcum)),
                 renk,
                 (_degiskenler(alt_tr, "tr", surum, olcum),
                  _degiskenler(alt_en, "en", surum, olcum)))
                for buyuk, renk, (alt_tr, alt_en) in kopya["kutular"]
            ]
        if "kartlar" in kopya:
            kopya["kartlar"] = [
                (basliklar,
                 (_degiskenler(govde_tr, "tr", surum, olcum),
                  _degiskenler(govde_en, "en", surum, olcum)))
                for basliklar, (govde_tr, govde_en) in kopya["kartlar"]
            ]
        cozulmus.append(kopya)
    return cozulmus


def sunum_kur(dil, palet, surum):
    """Tüm slaytları kurar → Presentation."""
    from pptx import Presentation
    from pptx.util import Inches

    sunum = Presentation()
    sunum.slide_width = Inches(SLAYT_G)
    sunum.slide_height = Inches(SLAYT_Y)

    for veri in _icerik(surum):
        if veri["tip"] == "kapak":
            _kapak_slaydi(sunum, palet, veri, dil, surum)
        elif veri["tip"] == "sayi":
            _sayi_slaydi(sunum, palet, veri, dil)
        elif veri["tip"] == "kart":
            _kart_slaydi(sunum, palet, veri, dil)
        elif veri["tip"] == "ekran":
            _ekran_slaydi(sunum, palet, veri, dil)
        else:
            raise ValueError("Bilinmeyen slayt tipi: %r" % veri["tip"])
    return sunum


# ══════════════════════════════════════════════════════════════════════════
# METİN SIĞMA DENETİMİ
#
# Kutular sabit; metin uzarsa PowerPoint onu küçültmez, TAŞIRIR ve slaytta
# üst üste biner. Bu yüzden üretim öncesi her metin bütçesine göre ölçülür.
# Denetim üretimi durdurmaz ama uyarır — bir cümlenin iki karakter aşması
# üretimi engellemeye değmez, sessiz kalması ise slaydı bozar.
# ══════════════════════════════════════════════════════════════════════════

def _olc(dil, surum):
    """Bütçeyi aşan metinleri döndürür: [(slayt_no, alan, uzunluk, bütçe, metin)]."""
    asanlar = []

    def bak(no, alan, metin):
        sinir = BUTCE[alan]
        if len(metin) > sinir:
            asanlar.append((no, alan, len(metin), sinir, metin))

    for no, veri in enumerate(_icerik(surum), 1):
        if veri["tip"] == "kapak":
            bak(no, "kapak_baslik", _M(veri["baslik"], dil))
            bak(no, "kapak_alt", _M(veri["alt"], dil))
            for pul in veri["puller"]:
                bak(no, "pul", _M(pul, dil))
            continue
        bak(no, "baslik", _M(veri["baslik"], dil))
        if veri["tip"] == "sayi":
            for _rakam, _rol, aciklama in veri["kutular"]:
                bak(no, "aciklama", _M(aciklama, dil))
            bak(no, "not", _M(veri["not"], dil))
        elif veri["tip"] == "ekran":
            bak(no, "alt", _M(veri["alt"], dil))
            for _dosya, etiket in veri["gorseller"]:
                if etiket:
                    bak(no, "ekran_etiket", _M(etiket, dil))
            if veri.get("dipnot"):
                bak(no, "dipnot", _M(veri["dipnot"], dil))
        else:
            bak(no, "alt", _M(veri["alt"], dil))
            for (_ikon, ad_tr, ad_en), govde in veri["kartlar"]:
                bak(no, "kart_ad", ad_tr if dil == "tr" else ad_en)
                bak(no, "kart_govde", _M(govde, dil))
            if veri.get("dipnot"):
                bak(no, "dipnot", _M(veri["dipnot"], dil))
    return asanlar


def kontrol(surum):
    """Üretmeden metin sığma denetimi yapar → 0 temiz."""
    toplam = 0
    for dil in DILLER:
        asanlar = _olc(dil, surum)
        etiket = dil.upper()
        if not asanlar:
            print("   %s [%s] %d slayt, metinlerin tamamı kutusuna sığıyor"
                  % (_iz("✓"), etiket, len(SLAYTLAR)))
            continue
        toplam += len(asanlar)
        print("   %s [%s] %d metin bütçeyi aşıyor:" % (_iz("✗"), etiket, len(asanlar)))
        for no, alan, uzunluk, sinir, metin in asanlar:
            print("      slayt %2d · %-12s %3d/%3d · %s…"
                  % (no, alan, uzunluk, sinir, metin[:60]))
    return 0 if toplam == 0 else 1


# ══════════════════════════════════════════════════════════════════════════
# PDF + PNG — PowerPoint COM ile dışa aktarma
#
# Neden LibreOffice değil de PowerPoint: sunum PowerPoint'te açılıp
# sunulacak. Başka bir dönüştürücü yazı tipi ölçüsünü ve emoji'yi farklı
# yorumlar; o zaman PDF ile ekranda görülen slayt aynı olmaz.
# ══════════════════════════════════════════════════════════════════════════

def _kilit_bekle(yol, saniye=15):
    """Dosya okunabilir hâle gelene kadar bekler.

    PowerPoint COM oturumu Close/Quit çağrıldıktan sonra da dosyayı kısa
    süre açık tutabiliyor; hemen ardından okumaya kalkan bir koşu
    "Invalid argument" ile düşer. Kilit saniyeler içinde kalkıyor, bu
    yüzden çökmek yerine beklemek doğru davranış.
    """
    import time
    son = None
    for _ in range(int(saniye * 4)):
        try:
            with open(yol, "rb") as f:
                f.read(4)
            return yol
        except OSError as hata:
            son = hata
            time.sleep(0.25)
    raise OSError("%s açılamadı (%s). PowerPoint hâlâ açık olabilir." % (yol, son))


def _com_oturumu(pptx, pdf, png_dizin):
    """Tek bir PowerPoint COM oturumu: aç → PDF yaz → (istenirse) PNG bas."""
    import time

    import win32com.client

    uygulama = win32com.client.Dispatch("PowerPoint.Application")
    sunum = None
    try:
        sunum = uygulama.Presentations.Open(os.path.abspath(pptx), ReadOnly=0,
                                            Untitled=0, WithWindow=0)
        sunum.SaveAs(os.path.abspath(pdf), 32)          # 32 = ppSaveAsPDF
        if png_dizin:
            sunum.Export(os.path.abspath(png_dizin), "PNG", 1920, 1080)
    finally:
        try:
            if sunum is not None:
                sunum.Close()
        except Exception:
            pass
        try:
            uygulama.Quit()
        except Exception:
            pass
        # Quit() hemen dönmez; süreç kapanırken açılan bir sonraki oturum
        # ÖLMEKTE OLAN örneği yakalar ve Presentations.Open "<unknown>.Open"
        # ile düşer (ölçüldü: TR'den hemen sonra EN üretimi bu yüzden çöktü).
        time.sleep(1.5)


def powerpoint_disa_aktar(pptx, pdf, png_dizin=None, png_uret=True, deneme=3):
    """PPTX'ten PDF ve (istenirse) slayt PNG'leri üretir → (pdf, [png]).

    COM çağrısı yeniden denenir: hata kalıcı değil, zamanlamaya bağlıdır.
    """
    import time

    if png_uret:
        png_dizin = png_dizin or os.path.join(tempfile.gettempdir(), "aky_sunum_png")
        if os.path.isdir(png_dizin):
            shutil.rmtree(png_dizin, ignore_errors=True)
        os.makedirs(png_dizin, exist_ok=True)

    for kalan in range(deneme, 0, -1):
        try:
            _com_oturumu(pptx, pdf, png_dizin if png_uret else None)
            break
        except Exception:
            if kalan == 1:
                raise
            time.sleep(3.0)

    if not png_uret:
        return pdf, []
    # Windows'ta glob büyük/küçük harf duyarsızdır: "*.PNG" ile "*.png" AYNI
    # dosyaları döndürür ve ikisini toplamak listeyi ikiye katlar.
    pngler = sorted({os.path.normcase(p) for p in glob.glob(os.path.join(png_dizin, "*.png"))},
                    key=lambda p: int(re.findall(r"(\d+)", os.path.basename(p))[-1]))
    return pdf, pngler


# ══════════════════════════════════════════════════════════════════════════
# HTML — tek dosyalık mobil sunum (slaytlar gömülü PNG)
# ══════════════════════════════════════════════════════════════════════════

_HTML_METIN = {
    "tr": {
        "kod": "tr",
        "baslik": "Akıllı Konaklama Yönetim Sistemi __SURUM__ — Tanıtım",
        "marka": "Akıllı Konaklama Yönetim Sistemi",
        "onceki": "Önceki slayt",
        "sonraki": "Sonraki slayt",
        "tamekran": "Tam ekran slayt",
        "kapat": "✕ Kapat",
        "ipucu": "Slayda dokun → tam ekran · telefonu yan çevirmek de olur",
        "uyari": "Yasal uyum sorumluluğu işletmeye aittir",
        "slayt": "Slayt",
    },
    "en": {
        "kod": "en",
        "baslik": "Smart Hospitality Management System __SURUM__ — Introduction",
        "marka": "Smart Hospitality Management System",
        "onceki": "Previous slide",
        "sonraki": "Next slide",
        "tamekran": "Full-screen slide",
        "kapat": "✕ Close",
        "ipucu": "Tap a slide → full screen · turning the phone sideways works too",
        "uyari": "Legal compliance remains the operator's responsibility",
        "slayt": "Slide",
    },
}

_HTML_KALIP = """<!doctype html>
<html lang="__DIL__">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1, viewport-fit=cover">
<meta name="theme-color" content="#0B0F14">
<title>__BASLIK__</title>
<style>
  *{box-sizing:border-box;margin:0;padding:0}
  html,body{height:100%;background:#0B0F14;color:#E6EDF3;
    font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,Calibri,sans-serif;
    overscroll-behavior:none;-webkit-text-size-adjust:100%}
  body{display:flex;flex-direction:column;height:100dvh}
  header{flex:0 0 auto;display:flex;align-items:center;gap:.6rem;
    padding:calc(.5rem + env(safe-area-inset-top)) .9rem .5rem;
    border-bottom:1px solid #1c2430;background:#0B0F14}
  .brand{font-size:.86rem;font-weight:600;letter-spacing:.02em;white-space:nowrap;
    overflow:hidden;text-overflow:ellipsis}
  .ver{flex:0 0 auto;font-size:.7rem;font-weight:700;color:#0B0F14;background:#4A9EE0;
    border-radius:99px;padding:.15rem .5rem}
  .count{margin-left:auto;flex:0 0 auto;font-size:.8rem;color:#8B98A9;font-variant-numeric:tabular-nums}
  .bar{flex:0 0 auto;height:2px;background:#1c2430}
  .bar > i{display:block;height:100%;width:0;background:#4A9EE0;transition:width .18s ease}
  main{flex:1 1 auto;display:flex;overflow-x:auto;overflow-y:hidden;
    scroll-snap-type:x mandatory;scroll-behavior:smooth;
    -webkit-overflow-scrolling:touch;scrollbar-width:none}
  main::-webkit-scrollbar{display:none}
  .slide{flex:0 0 100%;scroll-snap-align:center;scroll-snap-stop:always;
    display:flex;align-items:center;justify-content:center;padding:.5rem}
  .slide img{max-width:100%;max-height:100%;width:auto;height:auto;
    border-radius:10px;border:1px solid #1c2430;display:block;cursor:zoom-in}
  #zoom{position:fixed;inset:0;z-index:50;background:#0B0F14;display:none;
    align-items:center;justify-content:center;cursor:zoom-out;
    padding-bottom:calc(3.4rem + env(safe-area-inset-bottom))}
  #zoom.on{display:flex}
  #zoom img{display:block;max-width:100%;max-height:100%}
  @media (orientation:portrait){
    #zoom img{width:min(100dvh, calc(100dvw * 16 / 9));height:auto;
      max-width:none;max-height:none;transform:rotate(90deg)}
  }
  #zoom .close{position:absolute;top:calc(.6rem + env(safe-area-inset-top));
    right:.8rem;z-index:2;font-size:.8rem;color:#8B98A9;background:#141920;
    border:1px solid #2B3442;border-radius:9px;padding:.4rem .7rem}
  #zoom .znav{position:absolute;z-index:2;left:0;right:0;
    bottom:calc(.6rem + env(safe-area-inset-bottom));
    display:flex;align-items:center;justify-content:space-between;
    gap:.6rem;padding:0 .8rem;pointer-events:none}
  #zoom .znav button{pointer-events:auto}
  #zoom .zcount{font-size:.78rem;color:#8B98A9;background:#0B0F14cc;
    border-radius:99px;padding:.3rem .7rem;font-variant-numeric:tabular-nums}
  footer{flex:0 0 auto;display:flex;align-items:center;gap:.5rem;
    padding:.5rem .9rem calc(.5rem + env(safe-area-inset-bottom));
    border-top:1px solid #1c2430}
  button{flex:0 0 auto;background:#141920;color:#E6EDF3;border:1px solid #2B3442;
    border-radius:9px;padding:.55rem .95rem;font:inherit;font-size:.85rem;cursor:pointer;
    -webkit-tap-highlight-color:transparent}
  button:active{background:#1c2430}
  button:disabled{opacity:.35}
  .hint{flex:1 1 auto;text-align:center;font-size:.74rem;color:#8B98A9;
    overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
  @media (orientation:landscape) and (max-height:520px){
    header,footer{padding-top:.25rem;padding-bottom:.25rem}
    .brand,.count{font-size:.72rem}
    .slide{padding:.25rem}
  }
</style>
</head>
<body>

<header>
  <span class="ver">__SURUM__</span>
  <span class="brand">__MARKA__</span>
  <span class="count" id="count">1 / __ADET__</span>
</header>
<div class="bar"><i id="bar"></i></div>

<main id="deck">
__SLAYTLAR__
</main>

<footer>
  <button id="prev" aria-label="__ONCEKI__">◀</button>
  <span class="hint" id="hint">__IPUCU__</span>
  <button id="next" aria-label="__SONRAKI__">▶</button>
</footer>

<div id="zoom" role="dialog" aria-label="__TAMEKRAN__">
  <span class="close">__KAPAT__</span>
  <img id="zoomimg" alt="">
  <div class="znav">
    <button id="zprev" aria-label="__ONCEKI__">◀</button>
    <span class="zcount" id="zcount"></span>
    <button id="znext" aria-label="__SONRAKI__">▶</button>
  </div>
</div>

<script>
(function () {
  var deck = document.getElementById('deck'),
      slides = deck.querySelectorAll('.slide'),
      n = slides.length,
      count = document.getElementById('count'),
      bar = document.getElementById('bar'),
      prev = document.getElementById('prev'),
      next = document.getElementById('next'),
      hint = document.getElementById('hint'),
      cur = 0;

  function render() {
    count.textContent = (cur + 1) + ' / ' + n;
    bar.style.width = ((cur + 1) / n * 100) + '%';
    prev.disabled = cur === 0;
    next.disabled = cur === n - 1;
  }
  function go(i) {
    cur = Math.max(0, Math.min(n - 1, i));
    slides[cur].scrollIntoView({ behavior: 'smooth', inline: 'center', block: 'nearest' });
    render();
  }
  prev.onclick = function () { go(cur - 1); };
  next.onclick = function () { go(cur + 1); };

  var t;
  deck.addEventListener('scroll', function () {
    clearTimeout(t);
    t = setTimeout(function () {
      var i = Math.round(deck.scrollLeft / deck.clientWidth);
      if (i !== cur) { cur = i; render(); }
    }, 90);
  }, { passive: true });

  var zoom = document.getElementById('zoom'),
      zoomimg = document.getElementById('zoomimg'),
      zcount = document.getElementById('zcount'),
      zprev = document.getElementById('zprev'),
      znext = document.getElementById('znext');

  function openZoom(i) {
    cur = Math.max(0, Math.min(n - 1, i));
    zoomimg.src = slides[cur].querySelector('img').src;
    zoomimg.alt = '__SLAYT__ ' + (cur + 1);
    zcount.textContent = (cur + 1) + ' / ' + n;
    zprev.disabled = cur === 0;
    znext.disabled = cur === n - 1;
    zoom.classList.add('on');
    render();
  }
  function closeZoom() {
    zoom.classList.remove('on');
    zoomimg.removeAttribute('src');
    go(cur);
  }
  for (var i = 0; i < n; i++) {
    (function (k) {
      slides[k].querySelector('img').addEventListener('click', function () { openZoom(k); });
    })(i);
  }
  zoom.addEventListener('click', closeZoom);
  zprev.addEventListener('click', function (e) { e.stopPropagation(); openZoom(cur - 1); });
  znext.addEventListener('click', function (e) { e.stopPropagation(); openZoom(cur + 1); });

  document.addEventListener('keydown', function (e) {
    var zoomed = zoom.classList.contains('on');
    if (e.key === 'Escape') { closeZoom(); return; }
    if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') {
      e.preventDefault(); zoomed ? openZoom(cur + 1) : go(cur + 1);
    }
    if (e.key === 'ArrowLeft' || e.key === 'PageUp') {
      e.preventDefault(); zoomed ? openZoom(cur - 1) : go(cur - 1);
    }
    if (!zoomed && e.key === 'Home') go(0);
    if (!zoomed && e.key === 'End') go(n - 1);
  });

  setTimeout(function () { hint.textContent = '__UYARI__'; }, 8000);
  render();
})();
</script>
</body>
</html>
"""


def html_uret(pngler, yol, surum, dil):
    """Slayt PNG'lerini tek dosyalık mobil sunuma gömer → yazılan bayt sayısı."""
    M = _HTML_METIN[dil]
    parcalar = []
    for indis, png in enumerate(pngler, 1):
        with open(png, "rb") as f:
            veri = base64.b64encode(f.read()).decode("ascii")
        yukleme = "eager" if indis <= 2 else "lazy"
        parcalar.append(
            '<figure class="slide" id="s%d"><img src="data:image/png;base64,%s" '
            'alt="%s %d" loading="%s" decoding="async"></figure>'
            % (indis, veri, M["slayt"], indis, yukleme))

    html = _HTML_KALIP
    for anahtar, deger in (("__DIL__", M["kod"]), ("__BASLIK__", M["baslik"]),
                           ("__MARKA__", M["marka"]), ("__ONCEKI__", M["onceki"]),
                           ("__SONRAKI__", M["sonraki"]), ("__TAMEKRAN__", M["tamekran"]),
                           ("__KAPAT__", M["kapat"]), ("__IPUCU__", M["ipucu"]),
                           ("__UYARI__", M["uyari"]), ("__SLAYT__", M["slayt"]),
                           ("__SURUM__", "v" + surum), ("__ADET__", str(len(pngler))),
                           ("__SLAYTLAR__", "\n".join(parcalar))):
        html = html.replace(anahtar, deger)
    with open(yol, "w", encoding="utf-8") as f:
        f.write(html)
    return len(html.encode("utf-8"))


# ══════════════════════════════════════════════════════════════════════════
# ÜRETİM
# ══════════════════════════════════════════════════════════════════════════

def _dil_uret(dil, surum, sadece_pptx=False):
    """Tek dil için PPTX (+PDF+HTML+baskı) üretir → 0 başarılı."""
    pptx, pdf, html = _yollar(dil)
    etiket = dil.upper()

    asanlar = _olc(dil, surum)
    if asanlar:
        print("   %s [%s] %d metin bütçeyi aşıyor (--kontrol ile ayrıntı)"
              % (_iz("✗"), etiket, len(asanlar)))

    sunum = sunum_kur(dil, PALET_EKRAN, surum)
    sunum.save(pptx)
    print("   %s [%s] PPTX : %d slayt" % (_iz("✓"), etiket, len(sunum.slides)))

    if sadece_pptx:
        return 0

    # PNG klasörü dile göre ayrı: aynı klasöre iki dil basılırsa ikinci koşu
    # birincinin görüntülerini siler ve HTML yanlış dilin slaytlarını gömer.
    png_dizin = os.path.join(tempfile.gettempdir(), "aky_sunum_png_" + dil)
    try:
        pdf_yolu, pngler = powerpoint_disa_aktar(_kilit_bekle(pptx), pdf, png_dizin)
    except Exception as hata:
        print("   %s [%s] PowerPoint dışa aktarımı başarısız: %s" % (_iz("✗"), etiket, hata))
        print("   PPTX güncel; PDF/HTML elle dışa aktarılabilir.")
        return 1

    print("   %s [%s] PDF  : %.1f MB" % (_iz("✓"), etiket, os.path.getsize(pdf_yolu) / 1048576))

    boyut = html_uret(pngler, html, surum, dil)
    print("   %s [%s] HTML : %d slayt, %.1f MB" % (_iz("✓"), etiket, len(pngler), boyut / 1048576))

    # ── Baskı sürümü: beyaz zemin, mono yazıcı ──────────────────────────
    try:
        b_pptx, b_pdf = _baski_yollari(dil)
        baski = sunum_kur(dil, PALET_BASKI, surum)
        baski.save(b_pptx)
        powerpoint_disa_aktar(_kilit_bekle(b_pptx), b_pdf, png_dizin=None, png_uret=False)
        print("   %s [%s] BASKI: %.1f MB (beyaz zemin)"
              % (_iz("✓"), etiket, os.path.getsize(b_pdf) / 1048576))
    except Exception as hata:
        print("   %s [%s] baskı sürümü üretilemedi: %s" % (_iz("✗"), etiket, hata))
        return 1
    return 0


def durum():
    """Üretilmiş dosyaların durumunu yazar."""
    for dil in DILLER:
        etiket = dil.upper()
        for ad, yol in zip(("PPTX", "PDF ", "HTML"), _yollar(dil)):
            varsa = "%.1f MB" % (os.path.getsize(yol) / 1048576) if os.path.exists(yol) else "yok"
            print("   [%s] %s : %s" % (etiket, ad, varsa))
        for ad, yol in zip(("BASKI PPTX", "BASKI PDF "), _baski_yollari(dil)):
            varsa = "%.1f MB" % (os.path.getsize(yol) / 1048576) if os.path.exists(yol) else "yok"
            print("   [%s] %s : %s" % (etiket, ad, varsa))


def main(argv=None):
    argv = list(argv if argv is not None else sys.argv[1:])
    surum = _surum()
    print("=" * 74)
    print("SUNUM URETIMI - Akilli Konaklama Yonetim Sistemi v%s" % surum)
    print("=" * 74)

    if "--kontrol" in argv:
        kod = kontrol(surum)
        print()
        durum()
        return kod

    diller = DILLER
    if "--dil" in argv:
        secim = argv[argv.index("--dil") + 1].lower()
        if secim not in DILLER:
            print("   %s bilinmeyen dil: %s (tr veya en)" % (_iz("✗"), secim))
            return 1
        diller = (secim,)

    sadece_pptx = "--sadece-pptx" in argv
    for dil in diller:
        kod = _dil_uret(dil, surum, sadece_pptx)
        if kod:
            return kod

    print("\n" + "=" * 74)
    dosya = (1 if sadece_pptx else 5) * len(diller)
    print("SONUC: %d dosya guncel (%s) - %d slayt"
          % (dosya, ", ".join(d.upper() for d in diller), len(SLAYTLAR)))
    print("=" * 74)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
